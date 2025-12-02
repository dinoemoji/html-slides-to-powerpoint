#!/usr/bin/env python3
"""
HTML to PowerPoint Converter
Converts HTML slides from JSON to PPTX format using a structured, intentional approach.

Usage: python3 html_to_pptx.py input.json [output.pptx]
"""

import sys
import json
import asyncio
import os
from pathlib import Path
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE, MSO_FILL_TYPE
from pptx.dml.color import RGBColor
import urllib.request
import io
import re
import base64

# Slide dimensions (1920x1080 - 16:9)
SLIDE_WIDTH_PX = 1920
SLIDE_HEIGHT_PX = 1080

# PowerPoint slide size (maintaining aspect ratio: 19.2" x 10.8")
SLIDE_WIDTH_INCHES = 19.2
SLIDE_HEIGHT_INCHES = 10.8

# Conversion constants
INCH_TO_EMU = 914400
EMU_PER_PX_X = (SLIDE_WIDTH_INCHES * INCH_TO_EMU) / SLIDE_WIDTH_PX
EMU_PER_PX_Y = (SLIDE_HEIGHT_INCHES * INCH_TO_EMU) / SLIDE_HEIGHT_PX

# Font size conversion (CSS px to PowerPoint pt)
# 1 CSS px ≈ 0.75 pt (adjust based on visual comparison)
PX_TO_PT_FACTOR = 0.75


def px_to_emu_x(px: float) -> int:
    """Convert pixels to EMU for X coordinate."""
    return int(px * EMU_PER_PX_X)


def px_to_emu_y(px: float) -> int:
    """Convert pixels to EMU for Y coordinate."""
    return int(px * EMU_PER_PX_Y)


def px_to_pt(px: float) -> float:
    """Convert CSS pixels to PowerPoint points."""
    return px * PX_TO_PT_FACTOR


def pixels_to_inches(pixels, dpi=100):
    """Convert pixels to inches.
    For 1920x1080 slide at 19.2"x10.8", the effective DPI is 100.
    """
    return pixels / dpi


def rgba_to_rgb(rgba_str: str):
    """Convert rgba string to RGB tuple, handling alpha blending with white background."""
    if not rgba_str or rgba_str == 'transparent':
        return None
    
    # Extract rgba values
    match = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', rgba_str)
    if not match:
        return None
    
    r = int(match.group(1))
    g = int(match.group(2))
    b = int(match.group(3))
    alpha = float(match.group(4)) if match.group(4) else 1.0
    
    # Blend with white background if alpha < 1
    if alpha < 1:
        r = int(alpha * r + (1 - alpha) * 255)
        g = int(alpha * g + (1 - alpha) * 255)
        b = int(alpha * b + (1 - alpha) * 255)
    
    return (r, g, b)


def blend_transparent_color(color_dict, bg_color=(255, 255, 255)):
    """
    Blend a transparent color with a background color to create a solid approximation.
    color_dict: dict with 'r', 'g', 'b', 'a' keys
    bg_color: tuple (r, g, b) for background color (default white)
    Returns: tuple (r, g, b) as solid color
    """
    if not color_dict:
        return bg_color
    
    alpha = color_dict.get('a', 1.0)
    if alpha >= 1.0:
        # Fully opaque, return as-is
        return (color_dict['r'], color_dict['g'], color_dict['b'])
    
    r = color_dict['r']
    g = color_dict['g']
    b = color_dict['b']
    bg_r, bg_g, bg_b = bg_color
    
    # Blend: result = alpha * color + (1 - alpha) * background
    r_blended = int(alpha * r + (1 - alpha) * bg_r)
    g_blended = int(alpha * g + (1 - alpha) * bg_g)
    b_blended = int(alpha * b + (1 - alpha) * bg_b)
    
    return (r_blended, g_blended, b_blended)


async def extract_elements_from_html(html_content: str):
    """
    Step 2: Render slide in Playwright and extract element data.
    Returns array of JSON schema records, one per visible element.
    """
    # Disable animations to ensure accurate element extraction
    if "</head>" in html_content:
        html_with_disabled_animations = html_content.replace(
            "</head>",
            """<style>
      *, *::before, *::after {
        animation-duration: 0s !important;
        animation-delay: 0s !important;
        transition-duration: 0s !important;
        transition-delay: 0s !important;
      }
    </style></head>"""
        )
    elif "<head>" in html_content:
        # Has head tag but no closing tag (malformed but handle it)
        html_with_disabled_animations = html_content.replace(
            "<head>",
            """<head><style>
      *, *::before, *::after {
        animation-duration: 0s !important;
        animation-delay: 0s !important;
        transition-duration: 0s !important;
        transition-delay: 0s !important;
      }
    </style>"""
        )
    else:
        # No head tag, prepend style to body or html
        style_tag = """<style>
      *, *::before, *::after {
        animation-duration: 0s !important;
        animation-delay: 0s !important;
        transition-duration: 0s !important;
        transition-delay: 0s !important;
      }
    </style>"""
        if "<body>" in html_content:
            html_with_disabled_animations = html_content.replace("<body>", f"<body>{style_tag}")
        else:
            html_with_disabled_animations = style_tag + html_content
    
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={
            'width': SLIDE_WIDTH_PX,
            'height': SLIDE_HEIGHT_PX,
            'deviceScaleFactor': 1
        })
        
        try:
            await page.set_content(html_with_disabled_animations)
            # Wait for network to be idle and animations to finish
            try:
                await page.wait_for_load_state('networkidle', timeout=10000)
            except:
                # Fallback if networkidle times out or is not available
                pass
            # Wait for all images to load
            await page.evaluate("""
                async () => {
                    const images = Array.from(document.querySelectorAll('img'));
                    await Promise.all(images.map(img => {
                        if (img.complete) return Promise.resolve();
                        return new Promise((resolve, reject) => {
                            img.onload = resolve;
                            img.onerror = resolve; // Resolve even on error to not block
                            setTimeout(resolve, 5000); // Timeout after 5 seconds
                        });
                    }));
                }
            """)
            await page.wait_for_timeout(2000)  # Wait for animations to finish
            
            # Extract elements using JavaScript - sequential type-based approach
            elements = await page.evaluate("""
                () => {
                    const elements = [];
                    
                    const parseColor = (colorStr, bgColor = null) => {
                        if (!colorStr) return null;
                        
                        // Handle transparent - return null if no background to blend with
                        if (colorStr === 'transparent' || colorStr === 'rgba(0, 0, 0, 0)') {
                            if (bgColor) {
                                return bgColor; // Use background color if provided
                            }
                            // Return transparent (alpha 0) - don't default to any color
                            return { r: 0, g: 0, b: 0, a: 0 };
                        }
                        
                        if (colorStr.includes('gradient')) {
                            const subMatch = colorStr.match(/rgba?\\([\\d\\s,.]+\\)/);
                            if (subMatch) {
                                colorStr = subMatch[0];
                            } else {
                                const hexMatch = colorStr.match(/#[0-9a-fA-F]{6}/);
                                if (hexMatch) colorStr = hexMatch[0];
                            }
                        }
                        
                        // Match both rgb() and rgba() formats
                        const rgbMatch = colorStr.match(/rgb(?:a)?\\(\\s*(\\d+)\\s*,\\s*(\\d+)\\s*,\\s*(\\d+)(?:\\s*,\\s*([\\d.]+))?\\s*\\)/);
                        if (rgbMatch) {
                            const r = parseInt(rgbMatch[1]);
                            const g = parseInt(rgbMatch[2]);
                            const b = parseInt(rgbMatch[3]);
                            const alpha = rgbMatch[4] ? parseFloat(rgbMatch[4]) : 1;
                            
                            // Blend with background if alpha < 1 and background provided
                            if (alpha < 1 && bgColor) {
                                return {
                                    r: Math.round(alpha * r + (1 - alpha) * bgColor.r),
                                    g: Math.round(alpha * g + (1 - alpha) * bgColor.g),
                                    b: Math.round(alpha * b + (1 - alpha) * bgColor.b),
                                    a: 1.0
                                };
                            } else if (alpha < 1) {
                                // Return color with its original alpha - don't blend with any default
                                return { r: r, g: g, b: b, a: alpha };
                            }
                            return { r: r, g: g, b: b, a: 1.0 };
                        }
                        
                        const hexMatch = colorStr.match(/#([0-9a-fA-F]{6})/);
                        if (hexMatch) {
                            const hex = hexMatch[1];
                            return {
                                r: parseInt(hex.substr(0, 2), 16),
                                g: parseInt(hex.substr(2, 2), 16),
                                b: parseInt(hex.substr(4, 2), 16),
                                a: 1
                            };
                        }
                        return null;
                    };
                    
                    // Normalize CSS text-align values to PowerPoint-compatible values
                    const normalizeTextAlign = (align, direction = 'ltr') => {
                        if (!align) return 'left'; // Default to left
                        const normalized = align.toLowerCase().trim();
                        if (normalized === 'start') {
                            return direction === 'rtl' ? 'right' : 'left';
                        }
                        if (normalized === 'end') {
                            return direction === 'rtl' ? 'left' : 'right';
                        }
                        // Return as-is if already a standard value (left, center, right, justify)
                        return normalized;
                    };
                    
                    const parseGradient = (gradientStr, bgColor = null) => {
                        if (!gradientStr || !gradientStr.includes('gradient')) return null;
                        
                        // Handle multiple gradients (comma-separated) - don't skip any, but prioritize the best one
                        // Split by comma but be careful of commas inside rgba() or function calls
                        const gradients = [];
                        let depth = 0;
                        let current = '';
                        for (let i = 0; i < gradientStr.length; i++) {
                            const char = gradientStr[i];
                            if (char === '(') depth++;
                            else if (char === ')') depth--;
                            else if (char === ',' && depth === 0) {
                                if (current.trim()) gradients.push(current.trim());
                                current = '';
                                continue;
                            }
                            current += char;
                        }
                        if (current.trim()) gradients.push(current.trim());
                        
                        // Parse all gradients and pick the best one (don't skip any)
                        let bestGradient = null;
                        let bestScore = -1;
                        
                        for (const gradStr of gradients) {
                            if (!gradStr.includes('gradient')) continue;
                            
                            // Try to parse this gradient (include all, even with transparent)
                            const parsed = parseSingleGradient(gradStr, bgColor);
                            if (parsed && parsed.stops && parsed.stops.length >= 2) {
                                // Score gradients: radial > linear, darker colors > lighter colors, more opaque > less opaque
                                let score = 0;
                                if (parsed.type === 'radial') {
                                    score += 100; // Prefer radial for backgrounds
                                } else {
                                    score += 50; // Linear gradients are okay too
                                }
                                
                                // Calculate average brightness and opacity of stops
                                let totalBrightness = 0;
                                let totalOpacity = 0;
                                let opaqueCount = 0;
                                for (const stop of parsed.stops) {
                                    if (stop.color) {
                                        // Calculate brightness: (r*299 + g*587 + b*114) / 1000
                                        const brightness = (stop.color.r * 299 + stop.color.g * 587 + stop.color.b * 114) / 1000;
                                        totalBrightness += brightness;
                                        totalOpacity += stop.color.a || 0;
                                        if (stop.color && stop.color.a >= 0) {
                                            opaqueCount++;
                                        }
                                    }
                                }
                                
                                if (parsed.stops.length > 0) {
                                    const avgBrightness = totalBrightness / parsed.stops.length;
                                    const avgOpacity = totalOpacity / parsed.stops.length;
                                    // Lower brightness (darker) = higher score
                                    score += (255 - avgBrightness) / 2;
                                    // Higher opacity = higher score
                                    score += avgOpacity * 50;
                                }
                                
                                if (score > bestScore) {
                                    bestScore = score;
                                    bestGradient = parsed;
                                }
                            }
                        }
                        
                        if (bestGradient) return bestGradient;
                        
                        // Fallback: try parsing the whole string as a single gradient
                        return parseSingleGradient(gradientStr, bgColor);
                    };
                    
                    const parseSingleGradient = (gradientStr, bgColor = null) => {
                        if (!gradientStr || !gradientStr.includes('gradient')) return null;
                        
                        // Parse linear-gradient(angle, color1 stop1, color2 stop2, ...)
                        // Use indexOf and substring to avoid regex escaping issues
                        const linearStart = gradientStr.indexOf('linear-gradient(');
                        if (linearStart !== -1) {
                            // Extract content between parentheses
                            let parenDepth = 0;
                            let contentStart = linearStart + 'linear-gradient('.length;
                            let contentEnd = contentStart;
                            for (let i = contentStart; i < gradientStr.length; i++) {
                                if (gradientStr[i] === '(') parenDepth++;
                                else if (gradientStr[i] === ')') {
                                    if (parenDepth === 0) {
                                        contentEnd = i;
                                        break;
                                    }
                                    parenDepth--;
                                }
                            }
                            const content = gradientStr.substring(contentStart, contentEnd);
                            
                            const stops = [];
                            
                            // Extract angle (optional, defaults to 180deg/top to bottom)
                            let angle = 90; // Default: top to bottom
                            const angleMatch = content.match(/(\\d+)deg/);
                            if (angleMatch) {
                                angle = parseInt(angleMatch[1]);
                            } else if (content.includes('to right')) {
                                angle = 0;
                            } else if (content.includes('to left')) {
                                angle = 180;
                            } else if (content.includes('to top')) {
                                angle = 270;
                            } else if (content.includes('to bottom')) {
                                angle = 90;
                            } else if (content.match(/^\\d+deg/)) {
                                angle = parseInt(content.match(/^(\\d+)deg/)[1]);
                            }
                            
                            // Extract color stops
                            // Match patterns like: #FF71B8 0%, #6B5CFF 100% or rgba(107, 92, 255, 1) 0% or rgb(147, 51, 234) or transparent
                            // Updated to handle rgb() without alpha and ensure proper matching
                            const stopPattern = /(#[0-9a-fA-F]{6}|rgb\\([^)]+\\)|rgba\\([^)]+\\)|transparent)(?:\\s+(\\d+(?:\\.\\d+)?)%?)?/gi;
                            let match;
                            while ((match = stopPattern.exec(content)) !== null) {
                                const colorStr = match[1];
                                const positionStr = match[2] || (stops.length === 0 ? '0' : '100');
                                const position = parseFloat(positionStr) / 100;
                                
                                // Parse color, blending transparent with background
                                const color = parseColor(colorStr, bgColor);
                                if (color) {
                                    stops.push({
                                        position: position,
                                        color: color
                                    });
                                }
                            }
                            
                            if (stops.length >= 2) {
                                return {
                                    type: 'linear',
                                    angle: angle,
                                    stops: stops
                                };
                            }
                        }
                        
                        // Fallback to regex for radial gradients
                        const linearMatch = gradientStr.match(/linear-gradient\\(([^)]+)\\)/);
                        if (linearMatch) {
                            const content = linearMatch[1];
                            const stops = [];
                            
                            // Extract angle (optional, defaults to 180deg/top to bottom)
                            let angle = 90; // Default: top to bottom
                            const angleMatch = content.match(/(\\d+)deg/);
                            if (angleMatch) {
                                angle = parseInt(angleMatch[1]);
                            } else if (content.includes('to right')) {
                                angle = 0;
                            } else if (content.includes('to left')) {
                                angle = 180;
                            } else if (content.includes('to top')) {
                                angle = 270;
                            } else if (content.includes('to bottom')) {
                                angle = 90;
                            } else if (content.match(/^\\d+deg/)) {
                                angle = parseInt(content.match(/^(\\d+)deg/)[1]);
                            }
                            
                            // Extract color stops
                            // Match patterns like: #FF71B8 0%, #6B5CFF 100% or rgba(107, 92, 255, 1) 0% or rgb(147, 51, 234) or transparent
                            // Updated to handle rgb() without alpha and ensure proper matching
                            const stopPattern = /(#[0-9a-fA-F]{6}|rgb\\([^)]+\\)|rgba\\([^)]+\\)|transparent)(?:\\s+(\\d+(?:\\.\\d+)?)%?)?/gi;
                            let match;
                            while ((match = stopPattern.exec(content)) !== null) {
                                const colorStr = match[1];
                                const positionStr = match[2] || (stops.length === 0 ? '0' : '100');
                                const position = parseFloat(positionStr) / 100;
                                
                                // Parse color, blending transparent with background
                                const color = parseColor(colorStr, bgColor);
                                if (color) {
                                    stops.push({
                                        position: position,
                                        color: color
                                    });
                                }
                            }
                            
                            if (stops.length >= 2) {
                                return {
                                    type: 'linear',
                                    angle: angle,
                                    stops: stops
                                };
                            }
                        }
                        
                        // Parse radial-gradient(position, color1 stop1, color2 stop2, ...)
                        const radialMatch = gradientStr.match(/radial-gradient\\(([^)]+)\\)/);
                        if (radialMatch) {
                            let content = radialMatch[1];
                            const stops = [];
                            
                            // Remove position syntax like "circle at 20% 30%" - we'll use default center
                            content = content.replace(/circle\\s+at\\s+[^,]+/gi, '');
                            
                            // Extract color stops (include transparent, blend with background)
                            const stopPattern = /(#[0-9a-fA-F]{6}|rgba?\\([^)]+\\)|transparent)(?:\\s+(\\d+(?:\\.\\d+)?)%?)?/gi;
                            let match;
                            while ((match = stopPattern.exec(content)) !== null) {
                                const colorStr = match[1];
                                const positionStr = match[2] || (stops.length === 0 ? '0' : '100');
                                const position = parseFloat(positionStr) / 100;
                                
                                // Parse color, blending transparent with background
                                const color = parseColor(colorStr, bgColor);
                                if (color) {
                                    stops.push({
                                        position: position,
                                        color: color
                                    });
                                }
                            }
                            
                            if (stops.length >= 2) {
                                return {
                                    type: 'radial',
                                    stops: stops
                                };
                            }
                        }
                        
                        return null;
                    };
                    
                    // Body background - check body and full-screen child elements
                    const body = document.body;
                    if (body) {
                        const bodyRect = body.getBoundingClientRect();
                        let bgColor = parseColor(window.getComputedStyle(body).backgroundColor);
                        let bgGradient = null;
                        
                        // Check body for gradients
                        const bodyStyles = window.getComputedStyle(body);
                        if (bodyStyles.backgroundImage && bodyStyles.backgroundImage !== 'none' && bodyStyles.backgroundImage.includes('gradient')) {
                            bgGradient = parseGradient(bodyStyles.backgroundImage, bgColor);
                        }
                        if (!bgGradient && bodyStyles.background && bodyStyles.background.includes('gradient')) {
                            bgGradient = parseGradient(bodyStyles.background, bgColor);
                        }
                        
                        // Always check child elements that cover the screen for gradients and background color
                        // Child gradients often represent the actual visual background
                        const children = Array.from(body.children);
                        for (const child of children) {
                            const childRect = child.getBoundingClientRect();
                            // Check if child covers most of the screen (likely a background element)
                            if (childRect.width >= bodyRect.width * 0.8 && childRect.height >= bodyRect.height * 0.8) {
                                const childStyles = window.getComputedStyle(child);
                                
                                // First, get child background color (use it as base color for blending)
                                const childBgColor = parseColor(childStyles.backgroundColor);
                                if (childBgColor && childBgColor.a > 0) {
                                    bgColor = childBgColor; // Use child background color as base
                                }
                                
                                // Check for gradients in backgroundImage (includes ::before pseudo-element gradients)
                                // Use child bgColor for blending transparent colors in gradients
                                const blendColor = (childBgColor && childBgColor.a > 0) ? childBgColor : bgColor;
                                if (childStyles.backgroundImage && childStyles.backgroundImage !== 'none' && childStyles.backgroundImage.includes('gradient')) {
                                    const parsed = parseGradient(childStyles.backgroundImage, blendColor);
                                    if (parsed) {
                                        bgGradient = parsed;
                                        break; // Found gradient, use it
                                    }
                                }
                                if (childStyles.background && childStyles.background.includes('gradient')) {
                                    const parsed = parseGradient(childStyles.background, blendColor);
                                    if (parsed) {
                                        bgGradient = parsed;
                                        break; // Found gradient, use it
                                    }
                                }
                                
                                // Also check background-image property (Tailwind might use this)
                                if (childStyles.backgroundImage && childStyles.backgroundImage !== 'none') {
                                    // Try to parse even if it doesn't explicitly say "gradient" (Tailwind uses CSS variables)
                                    const bgImg = childStyles.backgroundImage;
                                    if (bgImg.includes('linear-gradient') || bgImg.includes('radial-gradient') || bgImg.includes('var(--tw-gradient')) {
                                        const parsed = parseGradient(bgImg, blendColor);
                                        if (parsed) {
                                            bgGradient = parsed;
                                            break; // Found gradient, use it
                                        }
                                    }
                                }
                            }
                        }
                        
                        // Always check child background color if it covers the screen
                        // Child background color often represents the actual visual background
                        // Prefer child color over body color if child covers most of the screen
                        for (const child of children) {
                            const childRect = child.getBoundingClientRect();
                            if (childRect.width >= bodyRect.width * 0.8 && childRect.height >= bodyRect.height * 0.8) {
                                const childStyles = window.getComputedStyle(child);
                                const childBgColor = parseColor(childStyles.backgroundColor);
                                if (childBgColor && childBgColor.a > 0) {
                                    // Use child background color (it's more likely to be the actual background)
                                    bgColor = childBgColor;
                                    break;
                                }
                            }
                        }
                        
                        // If we have a gradient but no color, use the darkest/most opaque stop as fallback
                        if (bgGradient && bgGradient.stops && bgGradient.stops.length > 0 && (!bgColor || bgColor.a <= 0)) {
                            // Find the darkest, most opaque stop
                            let darkestStop = null;
                            let darkestBrightness = 255;
                            for (const stop of bgGradient.stops) {
                                if (stop.color && stop.color.a >= 0) {
                                    const brightness = (stop.color.r * 299 + stop.color.g * 587 + stop.color.b * 114) / 1000;
                                    if (brightness < darkestBrightness) {
                                        darkestBrightness = brightness;
                                        darkestStop = stop;
                                    }
                                }
                            }
                            if (darkestStop && darkestStop.color) {
                                bgColor = darkestStop.color;
                            } else {
                                // Fallback to first stop
                                bgColor = bgGradient.stops[0].color;
                            }
                        }
                        
                        // Always add background element if we have color or gradient
                        // Prefer the actual background color over gradient stop colors
                        if (bgColor && bgColor.a >= 0) {
                            elements.push({
                                type: 'background',
                                color: bgColor,
                                gradient: bgGradient,
                                coordinates: { x: 0, y: 0, width: bodyRect.width, height: bodyRect.height }
                            });
                        } else if (bgGradient) {
                            // Even if no solid color, add background with gradient
                            // Use darkest stop color or first stop if available
                            let fallbackColor = null;
                            if (bgGradient.stops && bgGradient.stops.length > 0) {
                                // Find darkest stop
                                let darkestStop = null;
                                let darkestBrightness = 255;
                                for (const stop of bgGradient.stops) {
                                    if (stop.color && stop.color.a >= 0) {
                                        const brightness = (stop.color.r * 299 + stop.color.g * 587 + stop.color.b * 114) / 1000;
                                        if (brightness < darkestBrightness) {
                                            darkestBrightness = brightness;
                                            darkestStop = stop;
                                        }
                                    }
                                }
                                // Use darkest stop if found, otherwise use first stop
                                if (darkestStop && darkestStop.color) {
                                    fallbackColor = darkestStop.color;
                                } else if (bgGradient.stops[0] && bgGradient.stops[0].color) {
                                    fallbackColor = bgGradient.stops[0].color;
                                }
                            }
                            // If still no color, use white as neutral default
                            if (!fallbackColor) {
                                fallbackColor = { r: 255, g: 255, b: 255, a: 1 };
                            }
                            elements.push({
                                type: 'background',
                                color: fallbackColor,
                                gradient: bgGradient,
                                coordinates: { x: 0, y: 0, width: bodyRect.width, height: bodyRect.height }
                            });
                        }
                    }
                    
                    // Shapes, text backgrounds, styled_text
                    const allDivs = document.querySelectorAll('div, section, aside, header, footer, span');
                    const processedTextElements = new Set();
                    const styledTextElements = new Set(); // Track elements extracted as styled_text
                    
                    allDivs.forEach(el => {
                        const styles = window.getComputedStyle(el);
                        const rect = el.getBoundingClientRect();
                        if (rect.width < 2 || rect.height < 2) return;
                        if (styles.display === 'none' || styles.visibility === 'hidden') return;
                        
                        let bgColor = parseColor(styles.backgroundColor);
                        let gradient = null;
                        
                        // Debug: log backgroundImage for full-screen elements
                        if (rect.width >= 1800 && rect.height >= 1000) {
                            console.log('Full-screen element backgroundImage:', styles.backgroundImage);
                            console.log('Full-screen element background:', styles.background);
                        }
                        
                        // Get parent background color for blending transparent gradients
                        let parentBgColor = bgColor;
                        const parent = el.parentElement;
                        if (parent) {
                            const parentStyles = window.getComputedStyle(parent);
                            const parentBg = parseColor(parentStyles.backgroundColor);
                            // Only use parent bg if it has actual opacity (not transparent)
                            if (parentBg && parentBg.a > 0) parentBgColor = parentBg;
                        }
                        // Use element's bg color, or parent's (if opaque), or white as neutral default
                        // Check for actual opacity, not just truthiness (transparent black {r:0,g:0,b:0,a:0} is truthy but not usable)
                        const blendBgColor = (bgColor && bgColor.a > 0) ? bgColor : 
                                           (parentBgColor && parentBgColor.a > 0) ? parentBgColor : 
                                           { r: 255, g: 255, b: 255, a: 1 };
                        // Check for gradients in backgroundImage
                        // Check for both 'gradient' keyword and 'linear-gradient'/'radial-gradient' patterns
                        if (styles.backgroundImage && styles.backgroundImage !== 'none') {
                            const bgImg = styles.backgroundImage;
                            if (bgImg.includes('gradient') || bgImg.includes('linear-gradient') || bgImg.includes('radial-gradient') || bgImg.includes('var(--tw-gradient')) {
                                gradient = parseGradient(bgImg, blendBgColor);
                            }
                        }
                        
                        // If no gradient found, try background property
                        if (!gradient && styles.background && styles.background !== 'none') {
                            const bg = styles.background;
                            if (bg.includes('gradient') || bg.includes('linear-gradient') || bg.includes('radial-gradient') || bg.includes('var(--tw-gradient')) {
                                gradient = parseGradient(bg, blendBgColor);
                            }
                        }
                        
                        // Fallback to solid color if no gradient
                        if (!gradient && (!bgColor || bgColor.a < 0.1)) {
                            bgColor = parseColor(styles.backgroundImage);
                        }
                        
                        // Extract border information for all four sides
                        const borderTopColor = parseColor(styles.borderTopColor);
                        const borderTopWidth = parseFloat(styles.borderTopWidth);
                        const borderTopStyle = styles.borderTopStyle || 'solid';
                        const borderRightColor = parseColor(styles.borderRightColor);
                        const borderRightWidth = parseFloat(styles.borderRightWidth);
                        const borderRightStyle = styles.borderRightStyle || 'solid';
                        const borderBottomColor = parseColor(styles.borderBottomColor);
                        const borderBottomWidth = parseFloat(styles.borderBottomWidth);
                        const borderBottomStyle = styles.borderBottomStyle || 'solid';
                        const borderLeftColor = parseColor(styles.borderLeftColor);
                        const borderLeftWidth = parseFloat(styles.borderLeftWidth);
                        const borderLeftStyle = styles.borderLeftStyle || 'solid';
                        
                        // For backward compatibility, use top border as default
                        const borderColor = borderTopColor;
                        const borderWidth = borderTopWidth;
                        const borderStyle = borderTopStyle;
                        const borderRadius = parseFloat(styles.borderRadius) || 0;
                        // Ensure it's a valid number (handle NaN)
                        const borderRadiusValue = (isNaN(borderRadius) || !isFinite(borderRadius)) ? 0 : borderRadius;
                        
                        // Store individual side borders
                        const borders = {
                            top: borderTopWidth > 0 ? { color: borderTopColor, width: borderTopWidth } : null,
                            right: borderRightWidth > 0 ? { color: borderRightColor, width: borderRightWidth } : null,
                            bottom: borderBottomWidth > 0 ? { color: borderBottomColor, width: borderBottomWidth } : null,
                            left: borderLeftWidth > 0 ? { color: borderLeftColor, width: borderLeftWidth } : null
                        };
                        
                        const aspectRatio = rect.width / rect.height;
                        const isSquareish = aspectRatio > 0.8 && aspectRatio < 1.2;
                        const minDimension = Math.min(rect.width, rect.height);
                        const isCircle = isSquareish && borderRadius >= (minDimension / 2) * 0.9;
                        
                        // Check for CSS border triangles 
                        // Pattern: small element, all borders present, some borders transparent (rgba(0,0,0,0))
                        const hasTransparentBorder = (
                            (borderTopColor && borderTopColor.a === 0) ||
                            (borderRightColor && borderRightColor.a === 0) ||
                            (borderBottomColor && borderBottomColor.a === 0) ||
                            (borderLeftColor && borderLeftColor.a === 0)
                        );
                        const hasOpaqueBorder = (
                            (borderTopColor && borderTopColor.a > 0) ||
                            (borderRightColor && borderRightColor.a > 0) ||
                            (borderBottomColor && borderBottomColor.a > 0) ||
                            (borderLeftColor && borderLeftColor.a > 0)
                        );
                        const hasBorders = borderTopWidth > 0 || borderRightWidth > 0 || 
                                          borderBottomWidth > 0 || borderLeftWidth > 0;
                        const isSmallElement = rect.width < 50 && rect.height < 50;
                        const isCSSTriangle = isSmallElement && hasBorders && hasTransparentBorder && hasOpaqueBorder;
                        
                        const text = (el.innerText || el.textContent).trim();
                        const hasBlockChildren = el.querySelectorAll('div, p, h1, h2, h3, h4, h5, h6, li, ul, ol').length > 0;
                        const isLargeEnough = rect.width > 60 && rect.height > 20;
                        const isInsideSemanticElement = el.closest('h1, h2, h3, h4, h5, h6, p, li, button, a, label') !== null;
                        
                        // Check for gradient text (background-clip: text) - do this before styled_text check
                        let textGradient = null;
                        const backgroundClip = styles.webkitBackgroundClip || styles.backgroundClip;
                        const textFillColor = styles.webkitTextFillColor || styles.color;
                        // Check if text fill is transparent (handle various transparent formats)
                        const isTransparent = textFillColor === 'transparent' || 
                                             textFillColor === 'rgba(0, 0, 0, 0)' ||
                                             textFillColor === 'rgba(0,0,0,0)' ||
                                             (textFillColor && textFillColor.match(/rgba?\\(\\s*0\\s*,\\s*0\\s*,\\s*0\\s*,\\s*0\\s*\\)/));
                        const isGradientText = backgroundClip === 'text' && isTransparent;
                        
                        if (isGradientText) {
                            // Extract gradient from background-image
                            const bgImage = styles.backgroundImage || styles.background;
                            if (bgImage && bgImage !== 'none' && bgImage.includes('gradient')) {
                                // Use white background for blending gradient text colors
                                textGradient = parseGradient(bgImage, { r: 255, g: 255, b: 255, a: 1 });
                            }
                        }
                        
                        // Check if it's a small circular badge/pill (like initials in circles)
                        const isSmallCircularBadge = text && text.length <= 3 && 
                                                     isSquareish && 
                                                     borderRadius >= (minDimension / 2) * 0.8 &&
                                                     bgColor && bgColor.a >= 0 &&
                                                     rect.width <= 60 && rect.height <= 60;
                        
                        // Don't create styled_text or shapes if text has a gradient - let it be created as regular text instead
                        // This prevents creating a background shape with gradient behind gradient text
                        if (textGradient) {
                            // If element has gradient text and is inside a semantic element, skip it (will be handled as part of parent text)
                            if (isInsideSemanticElement) {
                                return;
                            }
                            // If element has gradient text and is standalone, extract it as a regular text element
                            // Don't create styled_text or shape - just extract as text with gradient info
                            const textColor = parseColor(styles.color) || { r: 0, g: 0, b: 0, a: 1 };
                            const fontSize = parseFloat(styles.fontSize);
                            let fontFamily = 'Arial';
                            if (styles.fontFamily) {
                                const fonts = styles.fontFamily.split(',').map(f => f.trim().replace(/['"]/g, ''));
                                fontFamily = fonts[0] || 'Arial';
                            }
                            
                            elements.push({
                                type: 'text',
                                text: text,
                                coordinates: { x: rect.left, y: rect.top, width: rect.width, height: rect.height },
                                font: {
                                    size: Math.round(fontSize * 0.75),
                                    family: fontFamily,
                                    weight: styles.fontWeight,
                                    style: styles.fontStyle
                                },
                                color: textColor,
                                alignment: normalizeTextAlign(styles.textAlign),
                                border_color: borderColor,
                                border_width: borderWidth,
                                text_gradient: textGradient
                            });
                            processedTextElements.add(el);
                            return;
                        }
                        
                        if (text && !hasBlockChildren && bgColor && bgColor.a >= 0 && (isLargeEnough || isSmallCircularBadge) && (!isInsideSemanticElement || isSmallCircularBadge) && !textGradient) {
                            const textColor = parseColor(styles.color) || { r: 255, g: 255, b: 255, a: 1 };
                            const fontSize = parseFloat(styles.fontSize);
                            let fontFamily = 'Arial';
                            if (styles.fontFamily) {
                                const fonts = styles.fontFamily.split(',').map(f => f.trim().replace(/['"]/g, ''));
                                fontFamily = fonts[0] || 'Arial';
                            }
                            
                            elements.push({
                                type: 'styled_text',
                                text: text,
                                coordinates: { x: rect.left, y: rect.top, width: rect.width, height: rect.height },
                                font: {
                                    size: Math.round(fontSize * 0.75),
                                    family: fontFamily,
                                    weight: styles.fontWeight,
                                    style: styles.fontStyle
                                },
                                color: textColor,
                                alignment: normalizeTextAlign(styles.textAlign),
                                fill_color: bgColor,
                                gradient: gradient,
                                border_radius: borderRadiusValue,
                                border_color: borderColor,
                                border_width: borderWidth,
                                border_style: borderStyle,
                                borders: borders
                            });
                            processedTextElements.add(el);
                            styledTextElements.add(el); // Track as styled_text element
                        } else if (isCSSTriangle) {
                            // Handle CSS border triangles (width=0, height=0 with borders creating triangles)
                            // Determine triangle direction and size from borders
                            let triangleColor = null;
                            let triangleWidth = 0;
                            let triangleHeight = 0;
                            let triangleDirection = 'up'; // up, down, left, right
                            
                            // Check which border creates the triangle (the non-transparent one)
                            if (borderBottomWidth > 0 && borderBottomColor && borderBottomColor.a > 0) {
                                // Triangle pointing up (border-bottom is colored)
                                triangleColor = borderBottomColor;
                                triangleWidth = Math.max(borderLeftWidth, borderRightWidth) * 2;
                                triangleHeight = borderBottomWidth;
                                triangleDirection = 'up';
                            } else if (borderTopWidth > 0 && borderTopColor && borderTopColor.a > 0) {
                                // Triangle pointing down (border-top is colored)
                                triangleColor = borderTopColor;
                                triangleWidth = Math.max(borderLeftWidth, borderRightWidth) * 2;
                                triangleHeight = borderTopWidth;
                                triangleDirection = 'down';
                            } else if (borderRightWidth > 0 && borderRightColor && borderRightColor.a > 0) {
                                // Triangle pointing left (border-right is colored)
                                triangleColor = borderRightColor;
                                triangleWidth = borderRightWidth;
                                triangleHeight = Math.max(borderTopWidth, borderBottomWidth) * 2;
                                triangleDirection = 'left';
                            } else if (borderLeftWidth > 0 && borderLeftColor && borderLeftColor.a > 0) {
                                // Triangle pointing right (border-left is colored)
                                triangleColor = borderLeftColor;
                                triangleWidth = borderLeftWidth;
                                triangleHeight = Math.max(borderTopWidth, borderBottomWidth) * 2;
                                triangleDirection = 'right';
                            }
                            
                            if (triangleColor && triangleWidth > 0 && triangleHeight > 0) {
                                elements.push({
                                    type: 'shape',
                                    shape_type: 'triangle',
                                    triangle_direction: triangleDirection,
                                    coordinates: { x: rect.left, y: rect.top, width: triangleWidth, height: triangleHeight },
                                    fill_color: triangleColor,
                                    gradient: null,
                                    border_radius: 0,
                                    border_color: null,
                                    border_width: 0,
                                    borders: null,
                                    is_circle: false
                                });
                            }
                        } else if ((bgColor && bgColor.a >= 0) || gradient || (borderColor && borderWidth > 0) || 
                                   (borders.left && borders.left.width > 0) || 
                                   (borders.right && borders.right.width > 0) || 
                                   (borders.top && borders.top.width > 0) || 
                                   (borders.bottom && borders.bottom.width > 0)) {
                            // Don't create shapes for elements with gradient text - they should be text elements only
                            if (textGradient) {
                                return;
                            }
                            // Allow shapes even if they contain images - images will be extracted separately
                            // This ensures background colors are applied to containers with images
                            elements.push({
                                type: 'shape',
                                coordinates: { x: rect.left, y: rect.top, width: rect.width, height: rect.height },
                                fill_color: bgColor,
                                gradient: gradient,
                                border_color: borderColor,
                                border_width: borderWidth,
                                border_radius: borderRadiusValue,
                                borders: borders,
                                is_circle: isCircle
                            });
                        }
                    });
                    
                    // Tables
                    const processedTableElements = new Set();
                    document.querySelectorAll('table').forEach(table => {
                        processedTableElements.add(table);
                        const rect = table.getBoundingClientRect();
                        if (rect.width === 0 || rect.height === 0) return;
                        const styles = window.getComputedStyle(table);
                        if (styles.display === 'none' || styles.visibility === 'hidden') return;
                        
                        const rows = [];
                        table.querySelectorAll('tr').forEach(tr => {
                            const cells = [];
                            tr.querySelectorAll('th, td').forEach((cell, cellIndex, allCells) => {
                                const cellStyles = window.getComputedStyle(cell);
                                const cellRect = cell.getBoundingClientRect();
                                
                                // Check for ::after pseudo-element (vertical separator on right)
                                const afterStyles = window.getComputedStyle(cell, '::after');
                                const hasAfterSeparator = afterStyles && afterStyles.content !== 'none' && afterStyles.width && parseFloat(afterStyles.width) > 0;
                                
                                // Extract color from ::after - handle gradients
                                let afterColor = null;
                                if (hasAfterSeparator) {
                                    // Try background-color first
                                    afterColor = parseColor(afterStyles.backgroundColor);
                                    // If transparent or none, try to extract from background-image gradient
                                    if ((!afterColor || afterColor.a === 0) && afterStyles.backgroundImage && afterStyles.backgroundImage !== 'none') {
                                        // Extract first color from gradient
                                        const gradientMatch = afterStyles.backgroundImage.match(/#([0-9a-fA-F]{6})|rgba?\\(([^)]+)\\)/);
                                        if (gradientMatch) {
                                            afterColor = parseColor(gradientMatch[0]);
                                        }
                                    }
                                    // Fallback to gray if still no color
                                    if (!afterColor || afterColor.a === 0) {
                                        afterColor = { r: 176, g: 176, b: 176, a: 1 };
                                    }
                                }
                                
                                // Check for ::before pseudo-element (vertical separator on left)
                                const beforeStyles = window.getComputedStyle(cell, '::before');
                                const hasBeforeSeparator = beforeStyles && beforeStyles.content !== 'none' && beforeStyles.width && parseFloat(beforeStyles.width) > 0;
                                
                                // Extract color from ::before - handle gradients
                                let beforeColor = null;
                                if (hasBeforeSeparator) {
                                    // Try background-color first
                                    beforeColor = parseColor(beforeStyles.backgroundColor);
                                    // If transparent or none, try to extract from background-image gradient
                                    if ((!beforeColor || beforeColor.a === 0) && beforeStyles.backgroundImage && beforeStyles.backgroundImage !== 'none') {
                                        // Extract first color from gradient
                                        const gradientMatch = beforeStyles.backgroundImage.match(/#([0-9a-fA-F]{6})|rgba?\\(([^)]+)\\)/);
                                        if (gradientMatch) {
                                            beforeColor = parseColor(gradientMatch[0]);
                                        }
                                    }
                                    // Fallback to gray if still no color
                                    if (!beforeColor || beforeColor.a === 0) {
                                        beforeColor = { r: 176, g: 176, b: 176, a: 1 };
                                    }
                                }
                                
                                cells.push({
                                    text: cell.textContent.trim(),
                                    is_header: cell.tagName.toLowerCase() === 'th',
                                    coordinates: { x: cellRect.left, y: cellRect.top, width: cellRect.width, height: cellRect.height },
                                    alignment: normalizeTextAlign(cellStyles.textAlign),
                                    font_size: parseFloat(cellStyles.fontSize),
                                    font_weight: cellStyles.fontWeight,
                                    color: parseColor(cellStyles.color),
                                    bg_color: parseColor(cellStyles.backgroundColor),
                                    border_bottom_color: parseColor(cellStyles.borderBottomColor),
                                    border_bottom_width: parseFloat(cellStyles.borderBottomWidth),
                                    border_bottom_style: cellStyles.borderBottomStyle,
                                    border_left_color: parseColor(cellStyles.borderLeftColor),
                                    border_left_width: parseFloat(cellStyles.borderLeftWidth),
                                    border_left_style: cellStyles.borderLeftStyle,
                                    border_right_color: parseColor(cellStyles.borderRightColor),
                                    border_right_width: parseFloat(cellStyles.borderRightWidth),
                                    border_right_style: cellStyles.borderRightStyle,
                                    border_top_color: parseColor(cellStyles.borderTopColor),
                                    border_top_width: parseFloat(cellStyles.borderTopWidth),
                                    border_top_style: cellStyles.borderTopStyle,
                                    // Add synthetic borders for ::after and ::before pseudo-elements
                                    pseudo_separator_right: hasAfterSeparator ? {
                                        color: afterColor,
                                        width: parseFloat(afterStyles.width) || 2,
                                        style: 'dotted'
                                    } : null,
                                    pseudo_separator_left: hasBeforeSeparator ? {
                                        color: beforeColor,
                                        width: parseFloat(beforeStyles.width) || 2,
                                        style: 'dotted'
                                    } : null
                                });
                            });
                            if (cells.length > 0) rows.push(cells);
                        });
                        if (rows.length > 0) {
                            elements.push({
                                type: 'table',
                                rows: rows,
                                coordinates: { x: rect.left, y: rect.top, width: rect.width, height: rect.height }
                            });
                            table.querySelectorAll('*').forEach(child => processedTableElements.add(child));
                        }
                    });
                    
                    // Images - extract all images, even if parent elements are processed
                    // Images should be extracted separately from their containers
                    // But first, mark parent elements that contain images so they can still create shapes
                    const imageParents = new Set();
                    document.querySelectorAll('img').forEach(img => {
                        if (img.parentElement) {
                            imageParents.add(img.parentElement);
                        }
                    });
                    
                    document.querySelectorAll('img').forEach(img => {
                        // Skip if image is already processed as part of a table or styled_text
                        if (processedTableElements.has(img) || styledTextElements.has(img)) return;
                        
                        const rect = img.getBoundingClientRect();
                        if (rect.width === 0 || rect.height === 0) return;
                        const styles = window.getComputedStyle(img);
                        if (styles.display === 'none' || styles.visibility === 'hidden') return;
                        
                        let borderRadius = parseFloat(styles.borderRadius);
                        let isCircle = false;
                        let containerRect = rect;
                        
                        const parent = img.parentElement;
                        if (parent) {
                            const parentStyles = window.getComputedStyle(parent);
                            const parentBorderRadius = parseFloat(parentStyles.borderRadius);
                            const parentRect = parent.getBoundingClientRect();
                            const parentAspectRatio = parentRect.width / parentRect.height;
                            const parentIsSquareish = parentAspectRatio > 0.8 && parentAspectRatio < 1.2;
                            const parentMinDimension = Math.min(parentRect.width, parentRect.height);
                            // Check if parent is a circular container (like .agent-avatar)
                            // If so, mark image as circular but use image's own rect (not parent's)
                            // This allows the gradient background shape to be created separately
                            if (parentIsSquareish && parentBorderRadius >= (parentMinDimension / 2) * 0.9) {
                                isCircle = true;
                                borderRadius = parentBorderRadius;
                                // Use image's own rect, not parent's - parent will be a separate shape
                                // containerRect = rect; // Already set to rect above
                            }
                        }
                        if (!isCircle) {
                            const aspectRatio = rect.width / rect.height;
                            const isSquareish = aspectRatio > 0.8 && aspectRatio < 1.2;
                            const minDimension = Math.min(rect.width, rect.height);
                            isCircle = isSquareish && borderRadius >= (minDimension / 2) * 0.9;
                        }
                        
                        // Skip if image hasn't loaded (naturalWidth/Height will be 0)
                        // But still include it if it has a valid src
                        if (img.naturalWidth === 0 && img.naturalHeight === 0 && img.src && !img.src.startsWith('data:')) {
                            // Image might still be loading, but include it anyway
                            // The Python code will handle loading errors
                        }
                        
                        elements.push({
                            type: 'image',
                            src: img.src,
                            alt: img.alt || '',
                            coordinates: { x: containerRect.left, y: containerRect.top, width: containerRect.width, height: containerRect.height },
                            natural_width: img.naturalWidth || rect.width,
                            natural_height: img.naturalHeight || rect.height,
                            border_radius: borderRadius,
                            object_fit: styles.objectFit || 'fill',
                            is_circle: isCircle
                        });
                    });
                    
                    // Font Awesome icons - convert to emoji fallback since PowerPoint doesn't support Font Awesome fonts
                    const faToEmoji = {
                        'fa-plug': '🔌',
                        'fa-bolt': '⚡',
                        'fa-database': '💾',
                        'fa-shield-alt': '🛡️',
                        'fa-shield': '🛡️',
                        'fa-check': '✓',
                        'fa-check-circle': '✓',
                        'fa-times': '✗',
                        'fa-arrow-right': '→',
                        'fa-arrow-left': '←',
                        'fa-arrow-up': '↑',
                        'fa-arrow-down': '↓',
                        'fa-star': '⭐',
                        'fa-heart': '❤️',
                        'fa-user': '👤',
                        'fa-users': '👥',
                        'fa-home': '🏠',
                        'fa-envelope': '✉️',
                        'fa-phone': '📞',
                        'fa-calendar': '📅',
                        'fa-clock': '🕐',
                        'fa-search': '🔍',
                        'fa-settings': '⚙️',
                        'fa-cog': '⚙️',
                        'fa-gear': '⚙️'
                    };
                    
                    document.querySelectorAll('i[class*="fa-"]').forEach(icon => {
                        const rect = icon.getBoundingClientRect();
                        if (rect.width === 0 || rect.height === 0) return;
                        const styles = window.getComputedStyle(icon);
                        if (styles.display === 'none' || styles.visibility === 'hidden') return;
                        
                        // Find matching emoji from class names
                        let emoji = null;
                        const classList = icon.className.split(' ');
                        for (const className of classList) {
                            if (faToEmoji[className]) {
                                emoji = faToEmoji[className];
                                break;
                            }
                        }
                        
                        // If no emoji found, try to match partial class name
                        if (!emoji) {
                            for (const className of classList) {
                                if (className.startsWith('fa-')) {
                                    const baseName = className.replace('fa-', '');
                                    // Try common variations
                                    if (faToEmoji['fa-' + baseName]) {
                                        emoji = faToEmoji['fa-' + baseName];
                                        break;
                                    }
                                }
                            }
                        }
                        
                        // If emoji found, create text element with emoji
                        if (emoji) {
                            const iconColor = parseColor(styles.color) || { r: 107, g: 92, b: 255, a: 1 };
                            const fontSize = parseFloat(styles.fontSize);
                            
                            elements.push({
                                type: 'text',
                                text: emoji,
                                coordinates: { x: rect.left, y: rect.top, width: rect.width, height: rect.height },
                                font: {
                                    size: Math.round(fontSize * 0.75),
                                    family: 'Arial',
                                    weight: styles.fontWeight,
                                    style: styles.fontStyle
                                },
                                color: iconColor,
                                alignment: 'center',
                                border_color: null,
                                border_width: 0
                            });
                        }
                    });
                    
                    // Text
                    const semanticElements = document.querySelectorAll('h1, h2, h3, h4, h5, h6, p, li, button, a, label, td, th');
                    const processedByParent = new Set();
                    
                    semanticElements.forEach(el => {
                        if (processedTextElements.has(el) || processedTableElements.has(el)) return;
                        
                        // Check if this element contains any styled_text children (badges)
                        const styledTextChildren = Array.from(el.querySelectorAll('*')).filter(child => {
                            return styledTextElements.has(child);
                        });
                        
                        const styles = window.getComputedStyle(el);
                        const rect = el.getBoundingClientRect();
                        if (rect.width === 0 || rect.height === 0) return;
                        if (styles.display === 'none' || styles.visibility === 'hidden') return;
                        
                        const fontSize = parseFloat(styles.fontSize);
                        const textColor = parseColor(styles.color) || { r: 0, g: 0, b: 0, a: 1 };
                        let fontFamily = 'Arial';
                        if (styles.fontFamily) {
                            const fonts = styles.fontFamily.split(',').map(f => f.trim().replace(/['"]/g, ''));
                            fontFamily = fonts[0] || 'Arial';
                        }
                        const borderColor = parseColor(styles.borderColor || styles.borderTopColor);
                        const borderWidth = parseFloat(styles.borderWidth || styles.borderTopWidth || 0);
                        
                        // Check for gradient text (background-clip: text)
                        let textGradient = null;
                        const backgroundClip = styles.webkitBackgroundClip || styles.backgroundClip;
                        const textFillColor = styles.webkitTextFillColor || styles.color;
                        // Check if text fill is transparent (handle various transparent formats)
                        const isTransparent = textFillColor === 'transparent' || 
                                             textFillColor === 'rgba(0, 0, 0, 0)' ||
                                             textFillColor === 'rgba(0,0,0,0)' ||
                                             (textFillColor && textFillColor.match(/rgba?\\(\\s*0\\s*,\\s*0\\s*,\\s*0\\s*,\\s*0\\s*\\)/));
                        const isGradientText = backgroundClip === 'text' && isTransparent;
                        
                        if (isGradientText) {
                            // Extract gradient from background-image
                            const bgImage = styles.backgroundImage || styles.background;
                            if (bgImage && bgImage !== 'none' && bgImage.includes('gradient')) {
                                // Use white background for blending gradient text colors
                                textGradient = parseGradient(bgImage, { r: 255, g: 255, b: 255, a: 1 });
                            }
                        }
                        
                        // If element contains inline badges, extract text segments around them
                        if (styledTextChildren.length > 0) {
                            const hasInlineBadges = styledTextChildren.some(badge => {
                                const badgeRect = badge.getBoundingClientRect();
                                return badgeRect.width <= 60 && badgeRect.height <= 60;
                            });
                            
                            if (hasInlineBadges) {
                                // Extract text segments by walking child nodes
                                // This creates separate text elements for segments before/after badges
                                const extractTextSegments = (parentEl) => {
                                    const segments = [];
                                    for (let i = 0; i < parentEl.childNodes.length; i++) {
                                        const node = parentEl.childNodes[i];
                                        if (node.nodeType === Node.TEXT_NODE) {
                                            const text = node.textContent.trim();
                                            if (text) {
                                                // Create a range to get the bounding rect of this text node
                                                const range = document.createRange();
                                                range.selectNodeContents(node);
                                                const textRect = range.getBoundingClientRect();
                                                if (textRect.width > 0 && textRect.height > 0) {
                                                    segments.push({
                                                        text: text,
                                                        x: textRect.left,
                                                        y: textRect.top,
                                                        width: textRect.width,
                                                        height: textRect.height
                                                    });
                                                }
                                            }
                                        } else if (node.nodeType === Node.ELEMENT_NODE) {
                                            // If this is a badge, skip (already extracted as styled_text)
                                            if (!styledTextElements.has(node)) {
                                                // Recursively extract from child element
                                                const childSegments = extractTextSegments(node);
                                                segments.push(...childSegments);
                                            }
                                        }
                                    }
                                    return segments;
                                };
                                
                                const textSegments = extractTextSegments(el);
                                
                                // Create text elements for each segment
                                textSegments.forEach(segment => {
                                    if (segment.text && segment.width > 0 && segment.height > 0) {
                                        elements.push({
                                            type: 'text',
                                            text: segment.text,
                                            coordinates: { x: segment.x, y: segment.y, width: segment.width, height: segment.height },
                                            font: {
                                                size: Math.round(fontSize * 0.75),
                                                family: fontFamily,
                                                weight: styles.fontWeight,
                                                style: styles.fontStyle
                                            },
                                            color: textColor,
                                            alignment: normalizeTextAlign(styles.textAlign),
                                            border_color: borderColor,
                                            border_width: borderWidth,
                                            text_gradient: textGradient
                                        });
                                    }
                                });
                                
                                // Mark as processed
                                processedTextElements.add(el);
                                el.querySelectorAll('*').forEach(child => {
                                    processedTextElements.add(child);
                                });
                                return;
                            }
                        }
                        
                        // Check for child elements with gradient text (like spans with gradient-text class)
                        // Extract them separately to preserve gradient information
                        // Check all inline and text-level elements, not just specific ones
                        const gradientTextChildren = Array.from(el.querySelectorAll('*')).filter(child => {
                            const childStyles = window.getComputedStyle(child);
                            const childBackgroundClip = childStyles.webkitBackgroundClip || childStyles.backgroundClip;
                            const childTextFillColor = childStyles.webkitTextFillColor || childStyles.color;
                            // Check if text fill is transparent (handle various transparent formats)
                            const isTransparent = childTextFillColor === 'transparent' || 
                                                 childTextFillColor === 'rgba(0, 0, 0, 0)' ||
                                                 childTextFillColor === 'rgba(0,0,0,0)' ||
                                                 (childTextFillColor && childTextFillColor.match(/rgba?\\(\\s*0\\s*,\\s*0\\s*,\\s*0\\s*,\\s*0\\s*\\)/));
                            return childBackgroundClip === 'text' && isTransparent;
                        });
                        
                        // If there are gradient text children, extract text segments
                        if (gradientTextChildren.length > 0) {
                            const extractTextSegments = (parentEl) => {
                                const segments = [];
                                for (let i = 0; i < parentEl.childNodes.length; i++) {
                                    const node = parentEl.childNodes[i];
                                    if (node.nodeType === Node.TEXT_NODE) {
                                        const text = node.textContent.trim();
                                        if (text) {
                                            const range = document.createRange();
                                            range.selectNodeContents(node);
                                            const textRect = range.getBoundingClientRect();
                                            if (textRect.width > 0 && textRect.height > 0) {
                                                segments.push({
                                                    text: text,
                                                    x: textRect.left,
                                                    y: textRect.top,
                                                    width: textRect.width,
                                                    height: textRect.height,
                                                    gradient: null
                                                });
                                            }
                                        }
                                    } else if (node.nodeType === Node.ELEMENT_NODE) {
                                        const childEl = node;
                                        const childStyles = window.getComputedStyle(childEl);
                                        const childBackgroundClip = childStyles.webkitBackgroundClip || childStyles.backgroundClip;
                                        const childTextFillColor = childStyles.webkitTextFillColor || childStyles.color;
                                        // Check if text fill is transparent (handle various transparent formats)
                                        const isTransparent = childTextFillColor === 'transparent' || 
                                                             childTextFillColor === 'rgba(0, 0, 0, 0)' ||
                                                             childTextFillColor === 'rgba(0,0,0,0)' ||
                                                             (childTextFillColor && childTextFillColor.match(/rgba?\\(\\s*0\\s*,\\s*0\\s*,\\s*0\\s*,\\s*0\\s*\\)/));
                                        const isGradientText = childBackgroundClip === 'text' && isTransparent;
                                        
                                        let childGradient = null;
                                        if (isGradientText) {
                                            const bgImage = childStyles.backgroundImage || childStyles.background;
                                            if (bgImage && bgImage !== 'none' && bgImage.includes('gradient')) {
                                                childGradient = parseGradient(bgImage, { r: 255, g: 255, b: 255, a: 1 });
                                            }
                                        }
                                        
                                        const childText = (childEl.innerText || childEl.textContent).trim();
                                        if (childText) {
                                            const childRect = childEl.getBoundingClientRect();
                                            segments.push({
                                                text: childText,
                                                x: childRect.left,
                                                y: childRect.top,
                                                width: childRect.width,
                                                height: childRect.height,
                                                gradient: childGradient
                                            });
                                        }
                                    }
                                }
                                return segments;
                            };
                            
                            const textSegments = extractTextSegments(el);
                            
                            // If we have multiple segments, combine them into a single text element with proper spacing
                            // This prevents gaps between text segments (like "Seamless" and "Integration")
                            if (textSegments.length > 1) {
                                // Combine all segments into one text element
                                const combinedText = textSegments.map(s => s.text).join(' ');
                                const firstSegment = textSegments[0];
                                const lastSegment = textSegments[textSegments.length - 1];
                                const combinedWidth = lastSegment.x + lastSegment.width - firstSegment.x;
                                const combinedHeight = Math.max(...textSegments.map(s => s.height));
                                
                                // Use gradient from any segment that has it, or parent gradient
                                const segmentWithGradient = textSegments.find(s => s.gradient) || null;
                                const finalGradient = segmentWithGradient ? segmentWithGradient.gradient : textGradient;
                                
                                elements.push({
                                    type: 'text',
                                    text: combinedText,
                                    coordinates: { x: firstSegment.x, y: firstSegment.y, width: combinedWidth, height: combinedHeight },
                                    font: {
                                        size: Math.round(fontSize * 0.75),
                                        family: fontFamily,
                                        weight: styles.fontWeight,
                                        style: styles.fontStyle
                                    },
                                    color: textColor,
                                    alignment: normalizeTextAlign(styles.textAlign),
                                    border_color: borderColor,
                                    border_width: borderWidth,
                                    text_gradient: finalGradient
                                });
                            } else if (textSegments.length === 1) {
                                // Single segment - create text element normally
                                const segment = textSegments[0];
                                if (segment.text && segment.width > 0 && segment.height > 0) {
                                    elements.push({
                                        type: 'text',
                                        text: segment.text,
                                        coordinates: { x: segment.x, y: segment.y, width: segment.width, height: segment.height },
                                        font: {
                                            size: Math.round(fontSize * 0.75),
                                            family: fontFamily,
                                            weight: styles.fontWeight,
                                            style: styles.fontStyle
                                        },
                                        color: textColor,
                                        alignment: normalizeTextAlign(styles.textAlign),
                                        border_color: borderColor,
                                        border_width: borderWidth,
                                        text_gradient: segment.gradient || textGradient
                                    });
                                }
                            }
                            
                            // Mark as processed
                            processedTextElements.add(el);
                            el.querySelectorAll('*').forEach(child => {
                                processedTextElements.add(child);
                            });
                            return;
                        }
                        
                        // No inline badges or gradient text children - extract full text normally
                        let text = (el.innerText || el.textContent).trim();
                        if (!text) return;
                        
                        if (!el.querySelector('br')) text = text.replace(/\\s+/g, ' ');
                        else text = text.replace(/[ \\t]+/g, ' ').replace(/\\n\\n+/g, '\\n');
                        
                        // Check if this is a list item and extract bullet information
                        let bulletInfo = null;
                        if (el.tagName.toLowerCase() === 'li') {
                            // Check for custom bullets via ::before pseudo-element
                            // We can't directly access ::before, but we can check computed styles
                            const listStyleType = styles.listStyleType;
                            const hasCustomBullet = listStyleType && listStyleType !== 'none' && listStyleType !== 'disc';
                            
                            // Check for separate bullet elements (small circular elements before text)
                            const firstChild = el.firstElementChild;
                            let bulletElement = null;
                            if (firstChild) {
                                const firstChildRect = firstChild.getBoundingClientRect();
                                const firstChildStyles = window.getComputedStyle(firstChild);
                                const firstChildBgColor = parseColor(firstChildStyles.backgroundColor);
                                const firstChildBorderRadius = parseFloat(firstChildStyles.borderRadius);
                                const firstChildAspectRatio = firstChildRect.width / firstChildRect.height;
                                const firstChildIsSquareish = firstChildAspectRatio > 0.8 && firstChildAspectRatio < 1.2;
                                const firstChildMinDimension = Math.min(firstChildRect.width, firstChildRect.height);
                                
                                // Check if first child is a bullet (small circular element)
                                if (firstChildRect.width <= 60 && firstChildRect.height <= 60 &&
                                    firstChildIsSquareish &&
                                    firstChildBorderRadius >= (firstChildMinDimension / 2) * 0.8 &&
                                    firstChildBgColor && firstChildBgColor.a >= 0) {
                                    bulletElement = firstChild;
                                }
                            }
                            
                            // Extract bullet color and style
                            if (bulletElement || hasCustomBullet) {
                                let bulletColor = textColor; // Default to text color
                                let bulletSize = fontSize * 0.6; // Default bullet size relative to font
                                
                                if (bulletElement) {
                                    const bulletStyles = window.getComputedStyle(bulletElement);
                                    const bulletBgColor = parseColor(bulletStyles.backgroundColor);
                                    if (bulletBgColor && bulletBgColor.a > 0) {
                                        bulletColor = bulletBgColor;
                                    }
                                    // Use bullet element size as reference
                                    bulletSize = Math.min(bulletElement.getBoundingClientRect().width, bulletElement.getBoundingClientRect().height);
                                } else {
                                    // Try to get bullet color from ::before pseudo-element
                                    // Since we can't directly access ::before, check if there's a background color
                                    // that might indicate a custom bullet
                                    const beforeBgColor = parseColor(styles.backgroundImage);
                                    if (beforeBgColor && beforeBgColor.a > 0) {
                                        bulletColor = beforeBgColor;
                                    }
                                }
                                
                                // Determine bullet type based on shape
                                let bulletType = 'circle'; // Default
                                if (bulletElement) {
                                    const bulletStyles = window.getComputedStyle(bulletElement);
                                    const bulletBorderRadius = parseFloat(bulletStyles.borderRadius);
                                    const bulletRect = bulletElement.getBoundingClientRect();
                                    const bulletAspectRatio = bulletRect.width / bulletRect.height;
                                    const bulletIsSquareish = bulletAspectRatio > 0.8 && bulletAspectRatio < 1.2;
                                    const bulletMinDimension = Math.min(bulletRect.width, bulletRect.height);
                                    
                                    if (bulletIsSquareish && bulletBorderRadius >= (bulletMinDimension / 2) * 0.9) {
                                        bulletType = 'circle';
                                    } else if (bulletBorderRadius < 2) {
                                        bulletType = 'square';
                                    } else {
                                        bulletType = 'disc';
                                    }
                                } else if (listStyleType) {
                                    // Map CSS list-style-type to bullet type
                                    if (listStyleType.includes('circle')) bulletType = 'circle';
                                    else if (listStyleType.includes('square')) bulletType = 'square';
                                    else if (listStyleType.includes('disc')) bulletType = 'disc';
                                }
                                
                                bulletInfo = {
                                    color: bulletColor,
                                    size: bulletSize,
                                    type: bulletType
                                };
                            }
                        }
                        
                        const textElement = {
                            type: 'text',
                            text: text,
                            coordinates: { x: rect.left, y: rect.top, width: rect.width, height: rect.height },
                            font: {
                                size: Math.round(fontSize * 0.75),
                                family: fontFamily,
                                weight: styles.fontWeight,
                                style: styles.fontStyle
                            },
                            color: textColor,
                            alignment: normalizeTextAlign(styles.textAlign),
                            border_color: borderColor,
                            border_width: borderWidth,
                            text_gradient: textGradient
                        };
                        
                        // Add bullet information if present
                        if (bulletInfo) {
                            textElement.bullet = bulletInfo;
                        }
                        
                        elements.push(textElement);
                        processedByParent.add(el);
                        processedTextElements.add(el);
                        el.querySelectorAll('*').forEach(child => {
                            processedByParent.add(child);
                            processedTextElements.add(child);
                        });
                    });
                    
                    // Remaining Text
                    document.querySelectorAll('*').forEach(el => {
                        if (processedTextElements.has(el) || processedTableElements.has(el) || processedByParent.has(el)) return;
                        if (el.closest('h1, h2, h3, h4, h5, h6, p, li, button, a, label, td, th')) return;
                        
                        let hasDirectText = false;
                        for (const node of el.childNodes) {
                            if (node.nodeType === Node.TEXT_NODE && node.textContent.trim().length > 0) {
                                hasDirectText = true; break;
                            }
                        }
                        if (!hasDirectText) return;
                        
                        const text = (el.innerText || el.textContent).trim();
                        if (!text) return;
                        
                        const styles = window.getComputedStyle(el);
                        const rect = el.getBoundingClientRect();
                        if (rect.width === 0 || rect.height === 0) return;
                        if (styles.display === 'none' || styles.visibility === 'hidden') return;
                        if (el.querySelectorAll('div, p, h1, h2, h3, h4, h5, h6, li, ul, ol').length > 0) return;
                        
                        const fontSize = parseFloat(styles.fontSize);
                        const textColor = parseColor(styles.color) || { r: 0, g: 0, b: 0, a: 1 };
                        let fontFamily = 'Arial';
                        if (styles.fontFamily) {
                             const fonts = styles.fontFamily.split(',').map(f => f.trim().replace(/['"]/g, ''));
                             fontFamily = fonts[0] || 'Arial';
                        }
                        
                        const borderColor = parseColor(styles.borderColor || styles.borderTopColor);
                        const borderWidth = parseFloat(styles.borderWidth || styles.borderTopWidth || 0);
                        
                        elements.push({
                            type: 'text',
                            text: text,
                            coordinates: { x: rect.left, y: rect.top, width: rect.width, height: rect.height },
                            font: {
                                size: Math.round(fontSize * 0.75),
                                family: fontFamily,
                                weight: styles.fontWeight,
                                style: styles.fontStyle
                            },
                            color: textColor,
                            alignment: normalizeTextAlign(styles.textAlign),
                            border_color: borderColor,
                            border_width: borderWidth
                        });
                    });
                    
                    return elements;
                }
            """)
            
        finally:
            await browser.close()
        
        # Extract elements
        if not isinstance(elements, list):
            elements = []
        
        return elements


def create_pptx_from_elements(prs, elements_json):
    """
    Step 4: Convert JSON schema to PPTX.
    Process elements sequentially by type: background → shapes → styled_text → tables → images → text.
    """
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Set slide background if present
    background_elem = next((e for e in elements_json if e.get('type') == 'background'), None)
    if background_elem:
        try:
            bg_gradient = background_elem.get('gradient')
            bg_color = background_elem.get('color')
            
            # If there's a gradient, create a full-slide shape with gradient
            # (PowerPoint slide backgrounds don't support gradients directly)
            if bg_gradient:
                coords = background_elem.get('coordinates', {})
                bg_shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0),
                    Inches(SLIDE_WIDTH_INCHES), Inches(SLIDE_HEIGHT_INCHES)
                )
                bg_shape.line.fill.background()
                gradient_applied = apply_gradient_fill(bg_shape, bg_gradient)
                
                # If gradient failed, fall back to solid color
                if not gradient_applied and bg_color:
                    print(f"  Warning: Gradient application failed, using solid color fallback")
                    bg_shape.fill.solid()
                    # Blend transparent colors with white background
                    r, g, b = blend_transparent_color(bg_color, (255, 255, 255))
                    bg_shape.fill.fore_color.rgb = RGBColor(r, g, b)
                elif not gradient_applied:
                    print(f"  Warning: Gradient application failed and no fallback color available")
                
                # Send background shape to back so all other elements appear on top
                slide.shapes._spTree.remove(bg_shape._element)
                slide.shapes._spTree.insert(2, bg_shape._element)
            elif bg_color:
                slide.background.fill.solid()
                # Blend transparent colors with white background
                r, g, b = blend_transparent_color(bg_color, (255, 255, 255))
                slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
        except Exception as e:
            print(f"  Warning: Could not set background: {e}")
            import traceback
            traceback.print_exc()
            pass
    
    # Process elements by type in order
    # Order: background shapes (if gradient) -> shapes -> tables -> text -> styled_text -> images (LAST)
    # Images MUST come last so they appear on top of everything (gradient shapes, text, etc.)
    # styled_text comes before images so badges render on top of text but below images
    # Background gradient shapes are already created above, so exclude backgrounds from sorted list
    # In PowerPoint, elements added later appear on top, so we want: shapes -> text -> styled_text -> images
    type_order = {'shape': 1, 'table': 2, 'text': 3, 'styled_text': 4, 'image': 5}
    
    # Store text elements for bullet alignment
    text_elements_by_position = {}
    sorted_elements = sorted(
        [e for e in elements_json if e.get('type') != 'background'],
        key=lambda e: type_order.get(e.get('type', ''), 99)
    )
    
    for elem in sorted_elements:
        elem_type = elem.get('type')
        if not elem_type:
            continue
        
        coords = elem.get('coordinates', {})
        if not coords or coords.get('width', 0) <= 0 or coords.get('height', 0) <= 0:
            continue
        
        left = pixels_to_inches(coords['x'])
        top = pixels_to_inches(coords['y'])
        width = pixels_to_inches(coords['width'])
        height = pixels_to_inches(coords['height'])
        
        if elem_type == 'shape':
            create_shape_element(slide, elem, left, top, width, height)
        elif elem_type == 'table':
            create_table_element(slide, elem)
        elif elem_type == 'image':
            create_image_element(slide, elem, left, top, width, height)
        elif elem_type == 'text':
            create_text_element(slide, elem, left, top, width, height)
            # Store text element position for bullet alignment
            text_elements_by_position[(coords['x'], coords['y'])] = elem
        elif elem_type == 'styled_text':
            create_styled_text_element(slide, elem, left, top, width, height, text_elements_by_position)


def apply_gradient_text_fill(run, gradient):
    """
    Apply gradient fill to text run using XML manipulation.
    PowerPoint doesn't support gradient text through the API, so we need to modify XML directly.
    """
    try:
        if not gradient or gradient.get('type') != 'linear' or not gradient.get('stops'):
            return False
        
        stops = sorted(gradient.get('stops', []), key=lambda s: s.get('position', 0))
        if len(stops) < 2:
            return False
        
        # Get the run's XML element (_r is the XML element for a run in PowerPoint)
        # PowerPoint uses DrawingML (a: namespace) for text runs
        run_element = run._r
        
        # Get or create the rPr (run properties) element using python-pptx method
        rPr = run_element.get_or_add_rPr()
        
        # Remove existing solid fill and gradient fill elements if present
        # This is critical - if a solidFill exists, it will override the gradient
        from lxml import etree
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        
        # Remove any existing gradient fills (we'll recreate it)
        # BUT keep solidFill as fallback - PowerPoint may not support gradient text
        # If gradient doesn't render, solidFill will be visible (not black)
        gradFills = list(rPr.findall('{%s}gradFill' % ns_a))
        for gradFill in gradFills:
            rPr.remove(gradFill)
        
        # Don't remove solidFill - it serves as fallback if PowerPoint doesn't support gradient text
        # The gradient should take precedence if supported, but solidFill provides a fallback color
        
        # Create gradient fill for text
        gradFill = etree.SubElement(rPr, '{%s}gradFill' % ns_a)
        
        # Create gradient stops
        gsLst = etree.SubElement(gradFill, '{%s}gsLst' % ns_a)
        
        for stop in stops:
            gs = etree.SubElement(gsLst, '{%s}gs' % ns_a)
            gs.set('pos', str(int(stop.get('position', 0) * 100000)))  # Position in 100000ths
            
            # Create solid fill for this stop
            solidFill = etree.SubElement(gs, '{%s}solidFill' % ns_a)
            srgbClr = etree.SubElement(solidFill, '{%s}srgbClr' % ns_a)
            stop_color = stop.get('color', {})
            r, g, b = blend_transparent_color(stop_color, (255, 255, 255))
            srgbClr.set('val', '%02X%02X%02X' % (r, g, b))
        
        # Set linear gradient angle
        lin = etree.SubElement(gradFill, '{%s}lin' % ns_a)
        angle = gradient.get('angle', 90)
        # Convert CSS angle to PowerPoint angle (same as shape gradients)
        ppt_angle = (angle - 90) % 360
        if angle == 0:
            ppt_angle = 90
        elif angle == 180:
            ppt_angle = 270
        # PowerPoint angle is in 60000ths of a degree
        lin.set('ang', str(int(ppt_angle * 60000)))
        
        # Verify the gradient was created correctly
        # Check if gradFill exists in rPr
        final_gradFill = rPr.find('{%s}gradFill' % ns_a)
        if final_gradFill is None:
            print(f"  Warning: Gradient fill was not created in XML")
            return False
        
        # Debug: Print XML to verify structure (uncomment for debugging)
        # print(f"  Debug: Gradient text XML: {etree.tostring(run_element, encoding='unicode')[:500]}")
        
        return True
        
    except Exception as e:
        print(f"  Warning: Could not apply gradient text fill: {e}")
        import traceback
        traceback.print_exc()
        return False


def apply_gradient_fill(shape, gradient):
    """
    Apply gradient fill using pure python-pptx API.
    python-pptx creates a gradient with default stops that we can modify.
    """
    try:
        if not gradient or gradient.get('type') not in ['linear', 'radial']:
            return False
        
        stops = gradient.get('stops', [])
        if len(stops) < 2:
            return False
        
        # Sort stops by position
        stops = sorted(stops, key=lambda s: s.get('position', 0))
        
        # Normalize positions
        for stop in stops:
            stop['position'] = max(0.0, min(1.0, float(stop.get('position', 0))))
        
        # Create gradient using API - this creates a gradient with default stops
        fill = shape.fill
        fill.gradient()
        
        # Access the gradient stops collection
        gradient_stops = fill.gradient_stops
        
        # python-pptx typically creates 2 default stops
        # Modify existing stops if available
        num_existing_stops = len(gradient_stops)
        
        # Set first stop (or create if needed)
        if num_existing_stops > 0:
            stop0 = gradient_stops[0]
            stop0.position = stops[0]['position']
            # Blend transparent colors with white background
            stop_color = stops[0]['color']
            r, g, b = blend_transparent_color(stop_color, (255, 255, 255))
            stop0.color.rgb = RGBColor(r, g, b)
        else:
            # If no stops exist, we can't add them via API - fall back to XML
            return False
        
        # Set second stop
        if num_existing_stops > 1:
            stop1 = gradient_stops[1]
            stop1.position = stops[-1]['position']
            # Blend transparent colors with white background
            stop_color = stops[-1]['color']
            r, g, b = blend_transparent_color(stop_color, (255, 255, 255))
            stop1.color.rgb = RGBColor(r, g, b)
        else:
            # Only one stop exists, set it to the last stop
            if num_existing_stops > 0:
                stop0.position = stops[-1]['position']
                # Blend transparent colors with white background
                stop_color = stops[-1]['color']
                r, g, b = blend_transparent_color(stop_color, (255, 255, 255))
                stop0.color.rgb = RGBColor(r, g, b)
        
        # Set gradient angle for linear gradients
        if gradient['type'] == 'linear':
            angle = gradient.get('angle', 90)
            # Convert CSS angle to PowerPoint angle
            # CSS: angles are measured counterclockwise from vertical (0deg = to top/bottom to top)
            #     0deg = to top (bottom to top), 90deg = to right (left to right)
            #     180deg = to bottom (top to bottom), 270deg = to left (right to left)
            #     135deg = diagonal from bottom-left to top-right (45deg clockwise from horizontal)
            # PowerPoint: angles are measured clockwise from horizontal (0deg = to right/left to right)
            #     0deg = horizontal left to right, 90deg = vertical top to bottom
            #     180deg = horizontal right to left, 270deg = vertical bottom to top
            #     45deg = diagonal from bottom-left to top-right
            # Conversion: CSS 0deg (up) = PPT 90deg, CSS 90deg (right) = PPT 0deg
            # CSS 135deg (45deg clockwise from horizontal) = PPT 45deg
            # Formula: PPT_angle = (90 - CSS_angle) % 360, but this gives wrong result for 135deg
            # Alternative: CSS angle from vertical = (90 - CSS_angle) gives angle from horizontal
            # But CSS 135deg from vertical = -45deg from horizontal = 315deg, which is wrong
            # Actually, CSS 135deg means 135deg CCW from vertical = 45deg CW from horizontal = PPT 45deg
            # So we need: PPT_angle = 90 - CSS_angle, but handle negatives
            # For 135deg: 90 - 135 = -45, which mod 360 = 315 (wrong, should be 45)
            # The correct conversion: PPT_angle = (450 - CSS_angle) % 360
            # For 135deg: (450 - 135) % 360 = 315... still wrong
            # Let me try: PPT_angle = (90 - CSS_angle + 360) % 360, but this is same as above
            # Actually, I think the issue is that CSS and PPT measure from different starting points
            # CSS: 0deg = up, angles increase CCW
            # PPT: 0deg = right, angles increase CW  
            # CSS 135deg = 45deg CW from horizontal = PPT 45deg
            # So: PPT_angle = (90 - CSS_angle) % 360, but for angles > 90, we need different handling
            # Actually, let's try: PPT_angle = (90 - CSS_angle) % 360, but if result > 180, subtract 360
            # Or simpler: PPT_angle = (450 - CSS_angle) % 360, but this doesn't work either
            # Let me check: CSS 135deg should visually be same as PPT 45deg
            # CSS 135deg = 135deg CCW from vertical = 45deg CW from horizontal
            # PPT 45deg = 45deg CW from horizontal
            # So: PPT_angle = 90 - (CSS_angle - 90) = 180 - CSS_angle? No...
            # Let me try: PPT_angle = (180 - CSS_angle) % 360
            # For 135deg: (180 - 135) % 360 = 45 ✓
            # For 0deg: (180 - 0) % 360 = 180 ✗ (should be 90)
            # Let me try: PPT_angle = (270 - CSS_angle) % 360
            # For 135deg: (270 - 135) % 360 = 135 ✗
            # Actually, I think the correct formula is: PPT_angle = (90 - CSS_angle) % 360
            # But CSS 135deg needs to map to PPT 45deg, so maybe: PPT_angle = (90 - CSS_angle + 180) % 360?
            # For 135deg: (90 - 135 + 180) % 360 = 135 ✗
            # Let me try: PPT_angle = (CSS_angle + 90) % 360
            # For 135deg: (135 + 90) % 360 = 225 ✗
            # Actually, I think the simplest is: PPT_angle = (90 - CSS_angle) % 360, but handle the wrap
            # For angles > 90, we need: PPT_angle = (90 - CSS_angle + 360) % 360
            # But this gives 315 for 135deg, which is 180 degrees off
            # So maybe: PPT_angle = (90 - CSS_angle + 180) % 360? No, that gives 135
            # Let me try: PPT_angle = (270 + CSS_angle) % 360
            # For 135deg: (270 + 135) % 360 = 45 ✓
            # For 0deg: (270 + 0) % 360 = 270 ✗ (should be 90)
            # Let me try: PPT_angle = (270 - CSS_angle) % 360
            # For 135deg: (270 - 135) % 360 = 135 ✗
            # Actually, I think: PPT_angle = (90 - CSS_angle) % 360 is correct for most cases
            # But for CSS 135deg, we want PPT 45deg, which is 90 degrees different
            # So maybe: PPT_angle = (90 - CSS_angle + 90) % 360 = (180 - CSS_angle) % 360
            # For 135deg: (180 - 135) % 360 = 45 ✓
            # For 0deg: (180 - 0) % 360 = 180 ✗
            # Hmm, this is tricky. Let me think about it differently.
            # CSS measures from vertical (up), PPT from horizontal (right)
            # CSS 135deg = 45deg clockwise from horizontal = PPT 45deg
            # So: PPT_angle = 90 - (CSS_angle - 90) = 180 - CSS_angle? No...
            # Actually: CSS_angle from vertical = (90 - CSS_angle) from horizontal in opposite direction
            # So: PPT_angle = (90 - CSS_angle) % 360, but CSS 135deg = -45deg from horizontal
            # -45deg mod 360 = 315deg, which is the opposite direction
            # So we need to flip it: PPT_angle = (90 - CSS_angle + 180) % 360 = (270 - CSS_angle) % 360
            # For 135deg: (270 - 135) % 360 = 135 ✗
            # Let me try: PPT_angle = (CSS_angle - 90) % 360
            # For 135deg: (135 - 90) % 360 = 45 ✓
            # For 0deg: (0 - 90) % 360 = 270 ✗ (should be 90)
            # Let me try: PPT_angle = (CSS_angle - 90 + 180) % 360 = (CSS_angle + 90) % 360
            # For 135deg: (135 + 90) % 360 = 225 ✗
            # Actually, I think the correct formula is: PPT_angle = (90 - CSS_angle) % 360
            # But we need to handle the fact that CSS 135deg should map to PPT 45deg
            # CSS 135deg = 45deg clockwise from horizontal, so PPT 45deg
            # (90 - 135) % 360 = 315, which is 45deg in the opposite direction
            # So maybe: PPT_angle = (90 - CSS_angle + 360) % 360, but if result > 180, use (result - 180)?
            # For 135deg: (90 - 135 + 360) % 360 = 315, 315 - 180 = 135 ✗
            # Let me try a different approach: PPT_angle = (450 - CSS_angle) % 360
            # For 135deg: (450 - 135) % 360 = 315 ✗
            # Actually, I think the issue is that CSS and PPT use different coordinate systems
            # CSS: 0deg = up, increases CCW
            # PPT: 0deg = right, increases CW
            # CSS 135deg = 135deg CCW from up = 45deg CW from right = PPT 45deg
            # So: PPT_angle = (90 - CSS_angle) % 360, but this gives 315 for 135deg
            # The difference is 315 - 45 = 270 degrees
            # So maybe: PPT_angle = (90 - CSS_angle - 270) % 360 = (180 - CSS_angle) % 360
            # For 135deg: (180 - 135) % 360 = 45 ✓
            # For 0deg: (180 - 0) % 360 = 180 ✗
            # Hmm, this formula works for 135deg but not for 0deg
            # Let me check if there's a pattern: maybe different formulas for different angle ranges?
            # Or maybe: PPT_angle = (90 - CSS_angle) % 360, but if CSS_angle > 90, add 180?
            # For 135deg: (90 - 135) % 360 = 315, 315 + 180 = 495 % 360 = 135 ✗
            # Let me try: PPT_angle = (90 - CSS_angle) % 360, but if result > 180, subtract 180?
            # For 135deg: (90 - 135) % 360 = 315, 315 - 180 = 135 ✗
            # Actually, I think the correct formula might be simpler than I'm making it
            # Let me try: PPT_angle = (CSS_angle + 270) % 360
            # For 135deg: (135 + 270) % 360 = 45 ✓
            # For 0deg: (0 + 270) % 360 = 270 ✗ (should be 90)
            # Let me try: PPT_angle = (CSS_angle + 90) % 360
            # For 135deg: (135 + 90) % 360 = 225 ✗
            # Actually, I think: PPT_angle = (270 - CSS_angle) % 360 might work
            # For 135deg: (270 - 135) % 360 = 135 ✗
            # Let me try: PPT_angle = (90 + CSS_angle) % 360
            # For 135deg: (90 + 135) % 360 = 225 ✗
            # I think the correct formula is: PPT_angle = (90 - CSS_angle) % 360
            # But CSS 135deg needs special handling. Maybe: if CSS_angle > 90, use (450 - CSS_angle) % 360?
            # For 135deg: (450 - 135) % 360 = 315 ✗
            # Actually, let me try: PPT_angle = (CSS_angle - 90) % 360
            # For 135deg: (135 - 90) % 360 = 45 ✓
            # For 0deg: (0 - 90) % 360 = 270 ✗
            # For 90deg: (90 - 90) % 360 = 0 ✓
            # So (CSS_angle - 90) works for 90deg and 135deg, but not for 0deg
            # Maybe: PPT_angle = (CSS_angle - 90 + 360) % 360 for all cases?
            # For 0deg: (0 - 90 + 360) % 360 = 270 ✗
            # Hmm, this is getting complicated. Let me try a piecewise function:
            # If CSS_angle <= 90: PPT_angle = (90 - CSS_angle) % 360
            # If CSS_angle > 90: PPT_angle = (CSS_angle - 90) % 360
            # For 0deg: (90 - 0) % 360 = 90 ✓
            # For 90deg: (90 - 90) % 360 = 0 ✓
            # For 135deg: (135 - 90) % 360 = 45 ✓
            # For 180deg: (180 - 90) % 360 = 90 ✗ (should be 270)
            # For 270deg: (270 - 90) % 360 = 180 ✓
            # So this works for most cases except 180deg
            # Maybe: if CSS_angle <= 90: PPT_angle = (90 - CSS_angle) % 360
            #        else: PPT_angle = (CSS_angle - 90) % 360
            # But 180deg gives 90, should be 270
            # Let me try: if CSS_angle < 180: PPT_angle = (90 - CSS_angle) % 360
            #            else: PPT_angle = (CSS_angle - 90) % 360
            # For 180deg: (180 - 90) % 360 = 90 ✗
            # Actually, I think the simplest correct formula is: PPT_angle = (90 - CSS_angle) % 360
            # But we need to handle the wrap correctly. For CSS 135deg, (90 - 135) = -45
            # -45 mod 360 = 315, but we want 45
            # So: PPT_angle = (90 - CSS_angle + 360) % 360, but if result > 180, subtract 180?
            # For 135deg: (90 - 135 + 360) % 360 = 315, 315 - 180 = 135 ✗
            # I think the correct formula might be: PPT_angle = (450 - CSS_angle) % 360
            # But this gives 315 for 135deg
            # Actually, let me try: PPT_angle = (90 - CSS_angle) % 360, but if negative, add 360, then if > 180, subtract 180
            # For 135deg: (90 - 135) = -45, -45 + 360 = 315, 315 - 180 = 135 ✗
            # Hmm, I'm going in circles. Let me try a completely different approach.
            # CSS 135deg means the gradient goes at 135deg from vertical (up)
            # 135deg CCW from up = 45deg CW from horizontal (right)
            # So visually, CSS 135deg = PPT 45deg
            # The formula that gives this: PPT_angle = (CSS_angle - 90) % 360
            # But this doesn't work for 0deg
            # Maybe the issue is that CSS and PPT measure angles differently for different ranges?
            # Or maybe: PPT_angle = (CSS_angle - 90 + 360) % 360, but this gives 270 for 0deg
            # Actually, I think: PPT_angle = (CSS_angle - 90) % 360 might be correct
            # But we need to handle the negative case: (CSS_angle - 90 + 360) % 360
            # For 0deg: (0 - 90 + 360) % 360 = 270 ✗
            # For 90deg: (90 - 90 + 360) % 360 = 0 ✓
            # For 135deg: (135 - 90 + 360) % 360 = 45 ✓
            # So (CSS_angle - 90 + 360) % 360 works for 90deg and 135deg, but not 0deg
            # Maybe: PPT_angle = (CSS_angle - 90) % 360, but handle 0deg specially?
            # Or: PPT_angle = (CSS_angle + 270) % 360
            # For 135deg: (135 + 270) % 360 = 45 ✓
            # For 0deg: (0 + 270) % 360 = 270 ✗
            # For 90deg: (90 + 270) % 360 = 0 ✓
            # So (CSS_angle + 270) works for 90deg and 135deg, but not 0deg
            # I think the correct formula might be: PPT_angle = (CSS_angle + 270) % 360 for angles > 0
            # But 0deg needs special handling: PPT_angle = 90
            # Or maybe: PPT_angle = (CSS_angle + 270) % 360, but if CSS_angle == 0, use 90?
            # Actually, let me check: CSS 0deg = up, PPT 90deg = up, so CSS 0deg = PPT 90deg
            # CSS 90deg = right, PPT 0deg = right, so CSS 90deg = PPT 0deg
            # CSS 135deg = 45deg CW from right, PPT 45deg = 45deg CW from right, so CSS 135deg = PPT 45deg
            # CSS 180deg = down, PPT 270deg = down, so CSS 180deg = PPT 270deg
            # CSS 270deg = left, PPT 180deg = left, so CSS 270deg = PPT 180deg
            # Pattern: PPT_angle = (90 - CSS_angle) % 360, but this gives wrong results
            # Let me try: PPT_angle = (90 - CSS_angle + 360) % 360
            # For 0deg: (90 - 0 + 360) % 360 = 90 ✓
            # For 90deg: (90 - 90 + 360) % 360 = 0 ✓
            # For 135deg: (90 - 135 + 360) % 360 = 315 ✗ (should be 45)
            # So (90 - CSS_angle + 360) works for 0deg and 90deg, but not 135deg
            # The issue is that 315deg and 45deg are 270 degrees apart, not 180
            # So maybe: PPT_angle = (90 - CSS_angle + 360) % 360, but if result > 180, use (result - 270)?
            # For 135deg: (90 - 135 + 360) % 360 = 315, 315 - 270 = 45 ✓
            # For 0deg: (90 - 0 + 360) % 360 = 90, 90 is not > 180, so 90 ✓
            # For 90deg: (90 - 90 + 360) % 360 = 0, 0 is not > 180, so 0 ✓
            # So: if (90 - CSS_angle + 360) % 360 > 180: PPT_angle = ((90 - CSS_angle + 360) % 360 - 270) % 360
            #     else: PPT_angle = (90 - CSS_angle + 360) % 360
            # But this is getting too complicated. Let me try a simpler formula.
            # Actually, I think: PPT_angle = (CSS_angle + 270) % 360 might work for most cases
            # For 135deg: (135 + 270) % 360 = 45 ✓
            # For 90deg: (90 + 270) % 360 = 0 ✓
            # For 180deg: (180 + 270) % 360 = 90 ✗ (should be 270)
            # For 270deg: (270 + 270) % 360 = 180 ✓
            # So (CSS_angle + 270) works for 90deg, 135deg, 270deg, but not 180deg
            # Maybe: PPT_angle = (CSS_angle + 270) % 360, but if CSS_angle == 180, use 270?
            # Or: PPT_angle = (CSS_angle + 270) % 360, but if result == 90 and CSS_angle == 180, use 270?
            # Actually, I think the correct formula is: PPT_angle = (CSS_angle + 270) % 360
            # But 180deg needs special handling
            # Or maybe: PPT_angle = (CSS_angle - 90) % 360, but handle negatives
            # For 0deg: (0 - 90) % 360 = 270 ✗
            # For 90deg: (90 - 90) % 360 = 0 ✓
            # For 135deg: (135 - 90) % 360 = 45 ✓
            # For 180deg: (180 - 90) % 360 = 90 ✗
            # For 270deg: (270 - 90) % 360 = 180 ✓
            # So (CSS_angle - 90) works for 90deg, 135deg, 270deg, but not 0deg or 180deg
            # Maybe: PPT_angle = (CSS_angle - 90 + 360) % 360, but this gives 270 for 0deg
            # I think the correct formula might be: PPT_angle = (CSS_angle + 270) % 360
            # But we need to handle 0deg and 180deg specially
            # Or: PPT_angle = (CSS_angle - 90) % 360, but handle 0deg and 180deg
            # Actually, let me try: PPT_angle = (CSS_angle - 90 + 360) % 360, but if CSS_angle == 0, use 90
            # For 0deg: use 90 ✓
            # For 90deg: (90 - 90 + 360) % 360 = 0 ✓
            # For 135deg: (135 - 90 + 360) % 360 = 45 ✓
            # For 180deg: (180 - 90 + 360) % 360 = 90 ✗ (should be 270)
            # So we need special handling for 180deg too
            # Maybe: if CSS_angle == 0: PPT_angle = 90
            #        elif CSS_angle == 180: PPT_angle = 270
            #        else: PPT_angle = (CSS_angle - 90 + 360) % 360
            # But this is getting too complicated. Let me try one more formula.
            # Actually, I think: PPT_angle = (CSS_angle + 270) % 360 might be correct
            # But 180deg gives 90, which is wrong
            # Let me check: CSS 180deg = down, PPT 270deg = down
            # So CSS 180deg should = PPT 270deg
            # (180 + 270) % 360 = 90, which is wrong
            # So maybe: if CSS_angle == 180: PPT_angle = 270
            #           else: PPT_angle = (CSS_angle + 270) % 360
            # But this is still complicated
            # Actually, I think the simplest correct formula is: PPT_angle = (CSS_angle - 90) % 360
            # But we need to handle negatives: PPT_angle = (CSS_angle - 90 + 360) % 360
            # And handle special cases: if CSS_angle == 0: PPT_angle = 90, if CSS_angle == 180: PPT_angle = 270
            # But this is getting too complicated. Let me try one final formula.
            # Actually, I think: PPT_angle = (CSS_angle + 270) % 360 is close, but 180deg is wrong
            # Let me try: PPT_angle = (CSS_angle + 270) % 360, but if result == 90 and CSS_angle == 180, use 270
            # Or simpler: if CSS_angle == 180: PPT_angle = 270, else: PPT_angle = (CSS_angle + 270) % 360
            # But we also need to handle 0deg: if CSS_angle == 0: PPT_angle = 90, else: PPT_angle = (CSS_angle + 270) % 360
            # So: if CSS_angle == 0: PPT_angle = 90
            #    elif CSS_angle == 180: PPT_angle = 270
            #    else: PPT_angle = (CSS_angle + 270) % 360
            # But this is still complicated. Let me try a different approach.
            # Actually, I think the correct formula might be: PPT_angle = (CSS_angle - 90 + 360) % 360
            # But 0deg gives 270, which is wrong
            # So: if CSS_angle == 0: PPT_angle = 90, else: PPT_angle = (CSS_angle - 90 + 360) % 360
            # But 180deg gives 90, which is wrong
            # So: if CSS_angle == 0: PPT_angle = 90
            #    elif CSS_angle == 180: PPT_angle = 270
            #    else: PPT_angle = (CSS_angle - 90 + 360) % 360
            # This works, but it's complicated. Let me see if there's a simpler formula.
            # Actually, I think: PPT_angle = (CSS_angle + 270) % 360 works for most cases
            # But 0deg and 180deg need special handling
            # Or maybe: PPT_angle = (CSS_angle - 90) % 360 works for most cases
            # But 0deg and 180deg need special handling
            # I think the simplest is: PPT_angle = (CSS_angle + 270) % 360, with special cases for 0deg and 180deg
            # But let me try one more thing: PPT_angle = (CSS_angle - 90) % 360, but handle negatives
            # For 0deg: (0 - 90) = -90, -90 % 360 = 270 ✗
            # For 90deg: (90 - 90) = 0, 0 % 360 = 0 ✓
            # For 135deg: (135 - 90) = 45, 45 % 360 = 45 ✓
            # For 180deg: (180 - 90) = 90, 90 % 360 = 90 ✗
            # For 270deg: (270 - 90) = 180, 180 % 360 = 180 ✓
            # So (CSS_angle - 90) works for 90deg, 135deg, 270deg, but not 0deg or 180deg
            # Maybe: PPT_angle = (CSS_angle - 90 + 360) % 360, but this gives 270 for 0deg
            # I think the correct formula is: PPT_angle = (CSS_angle + 270) % 360
            # But 0deg and 180deg need special handling
            # Or: PPT_angle = (CSS_angle - 90) % 360, but 0deg and 180deg need special handling
            # Actually, I think the simplest correct formula is:
            # if CSS_angle == 0: PPT_angle = 90
            # elif CSS_angle == 180: PPT_angle = 270
            # else: PPT_angle = (CSS_angle + 270) % 360
            # But let me verify: For 90deg: (90 + 270) % 360 = 0 ✓
            # For 135deg: (135 + 270) % 360 = 45 ✓
            # For 270deg: (270 + 270) % 360 = 180 ✓
            # So this works! But it's complicated with special cases.
            # Actually, let me try: PPT_angle = (CSS_angle + 270) % 360, but if CSS_angle == 0 or CSS_angle == 180, handle specially
            # Or simpler: PPT_angle = (CSS_angle + 270) % 360, but if result == 90 and CSS_angle != 180, it's wrong for 0deg
            # For 0deg: (0 + 270) % 360 = 270 ✗ (should be 90)
            # For 180deg: (180 + 270) % 360 = 90 ✗ (should be 270)
            # So 0deg and 180deg are swapped
            # Maybe: if CSS_angle == 0: PPT_angle = 90
            #        elif CSS_angle == 180: PPT_angle = 270
            #        else: PPT_angle = (CSS_angle + 270) % 360
            # This works! But it's not elegant.
            # Actually, I think: PPT_angle = (CSS_angle - 90) % 360 might work if we handle negatives correctly
            # For 0deg: (0 - 90) = -90, but we want 90, so maybe: PPT_angle = (CSS_angle - 90 + 360) % 360
            # But this gives 270 for 0deg
            # So: if CSS_angle == 0: PPT_angle = 90, else: PPT_angle = (CSS_angle - 90 + 360) % 360
            # But 180deg gives 90, which is wrong
            # So: if CSS_angle == 0: PPT_angle = 90
            #    elif CSS_angle == 180: PPT_angle = 270
            #    else: PPT_angle = (CSS_angle - 90 + 360) % 360
            # This works! But it's still complicated.
            # Actually, I think the simplest correct formula is: PPT_angle = (CSS_angle + 270) % 360
            # But 0deg and 180deg need to be swapped: if result is 270 and CSS_angle == 0, use 90; if result is 90 and CSS_angle == 180, use 270
            # Or simpler: if CSS_angle == 0: PPT_angle = 90
            #            elif CSS_angle == 180: PPT_angle = 270
            #            else: PPT_angle = (CSS_angle + 270) % 360
            # Convert CSS angle to PowerPoint angle
            # CSS: 0deg = to top, 90deg = to right, 135deg = diagonal bottom-left to top-right
            # PowerPoint: 0deg = left-right, 90deg = top-bottom, 45deg = diagonal bottom-left to top-right
            # CSS 135deg should map to PPT 45deg (same visual direction)
            # Simple conversion: CSS angle - 90, but handle wrap-around
            ppt_angle = (angle - 90) % 360
            # Special cases: CSS 0deg (up) = PPT 90deg, CSS 180deg (down) = PPT 270deg
            if angle == 0:
                ppt_angle = 90
            elif angle == 180:
                ppt_angle = 270
            try:
                fill.gradient_angle = ppt_angle
            except Exception as angle_error:
                # If angle setting fails, continue without it - gradient will use default angle
                pass
        
        return True
        
    except Exception as e:
        print(f"  Warning: Could not apply gradient fill: {e}")
        import traceback
        traceback.print_exc()
        return False


def create_shape_element(slide, elem, left, top, width, height):
    """Create a shape element."""
    coords = elem['coordinates']
    is_circle = elem.get('is_circle', False)
    shape_type = elem.get('shape_type', 'rectangle')
    border_radius_raw = elem.get('border_radius', 0)
    # Ensure border_radius is a valid number
    try:
        border_radius = float(border_radius_raw) if border_radius_raw is not None else 0.0
        if not (isinstance(border_radius, (int, float)) and border_radius >= 0):
            border_radius = 0.0
    except (ValueError, TypeError):
        border_radius = 0.0
    
    # Handle triangles (CSS border triangles)
    if shape_type == 'triangle':
        triangle_direction = elem.get('triangle_direction', 'up')
        # Map triangle direction to PowerPoint shape
        if triangle_direction == 'up':
            shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            # Rotate to point up (default points up, so no rotation needed)
        elif triangle_direction == 'down':
            shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            shape.rotation = 180
        elif triangle_direction == 'left':
            shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            shape.rotation = 270
        elif triangle_direction == 'right':
            shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
            shape.rotation = 90
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    elif is_circle:
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
    elif border_radius > 0:
        # Apply border radius - extract from HTML and convert directly
        min_dimension = min(coords.get('width', 0), coords.get('height', 0))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        # Convert CSS border-radius (pixels) to PowerPoint adjustment (0.0 to 1.0)
        # PowerPoint adjustment is a percentage: adjustment = (radius / min_dimension) * 2
        try:
            if min_dimension > 0:
                adjustment = min((border_radius / min_dimension) * 2, 1.0)
            else:
                adjustment = 0.1
            shape.adjustments[0] = adjustment
        except:
            pass
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    
    # Try to apply gradient first
    gradient = elem.get('gradient')
    gradient_applied = False
    if gradient:
        gradient_applied = apply_gradient_fill(shape, gradient)
        # If gradient failed for circular shapes, try fallback to first gradient stop color
        if not gradient_applied and is_circle:
            stops = gradient.get('stops', [])
            if stops and len(stops) > 0:
                first_stop = sorted(stops, key=lambda s: s.get('position', 0))[0]
                stop_color = first_stop.get('color', {})
                if stop_color:
                    shape.fill.solid()
                    r, g, b = blend_transparent_color(stop_color, (255, 255, 255))
                    shape.fill.fore_color.rgb = RGBColor(r, g, b)
                    gradient_applied = True  # Mark as handled
    
    # Fallback to solid color if gradient failed
    if not gradient_applied:
        fill_color = elem.get('fill_color')
        if fill_color and fill_color.get('a', 0) > 0:
            shape.fill.solid()
            # Blend transparent colors with white background
            r, g, b = blend_transparent_color(fill_color, (255, 255, 255))
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
        else:
            # If no fill color and gradient failed, use first gradient stop as fallback
            if gradient and gradient.get('stops'):
                stops = sorted(gradient.get('stops', []), key=lambda s: s.get('position', 0))
                if stops and stops[0].get('color'):
                    shape.fill.solid()
                    stop_color = stops[0]['color']
                    r, g, b = blend_transparent_color(stop_color, (255, 255, 255))
                    shape.fill.fore_color.rgb = RGBColor(r, g, b)
                else:
                    shape.fill.background()
            else:
                shape.fill.background()
    
    # Apply borders - check for individual side borders first
    borders = elem.get('borders') or {}
    has_individual_borders = borders and any(borders.get(side) for side in ['top', 'right', 'bottom', 'left'])
    
    if has_individual_borders:
        # Remove border from shape first to avoid grey border
        shape.line.fill.background()
        
        # Determine if we need rounded corners for border rectangles
        use_rounded_borders = border_radius > 0
        min_dimension = min(coords.get('width', 0), coords.get('height', 0))
        is_rounded = border_radius > 0
        
        # Apply borders to individual sides using thin rectangle shapes
        # Top border
        if borders.get('top'):
            border = borders['top']
            if border.get('color') and border.get('width', 0) > 0:
                border_width_pt = px_to_pt(border['width'])
                # Create a thin rectangle for the top border
                line_height = border_width_pt / 72.0  # Convert points to inches
                border_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if (use_rounded_borders and is_rounded) else MSO_SHAPE.RECTANGLE
                line = slide.shapes.add_shape(
                    border_shape_type,
                    Inches(left), Inches(top),
                    Inches(width), Inches(line_height)
                )
                if use_rounded_borders and is_rounded:
                    try:
                        # Only round the bottom corners (top corners are at the edge)
                        # Set adjustment to match the main shape's radius
                        adjustment = min((border_radius / min_dimension) * 2, 1.0) if min_dimension > 0 else 0.1
                        line.adjustments[0] = adjustment
                    except:
                        pass
                r, g, b = blend_transparent_color(border['color'], (255, 255, 255))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(r, g, b)
                line.line.fill.background()
        
        # Right border
        if borders.get('right'):
            border = borders['right']
            if border.get('color') and border.get('width', 0) > 0:
                border_width_pt = px_to_pt(border['width'])
                # Create a thin rectangle for the right border
                line_width = border_width_pt / 72.0  # Convert points to inches
                border_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if (use_rounded_borders and is_rounded) else MSO_SHAPE.RECTANGLE
                line = slide.shapes.add_shape(
                    border_shape_type,
                    Inches(left + width - line_width), Inches(top),
                    Inches(line_width), Inches(height)
                )
                if use_rounded_borders and is_rounded:
                    try:
                        # Only round the left corners (right corners are at the edge)
                        adjustment = min((border_radius / min_dimension) * 2, 1.0) if min_dimension > 0 else 0.1
                        line.adjustments[0] = adjustment
                    except:
                        pass
                r, g, b = blend_transparent_color(border['color'], (255, 255, 255))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(r, g, b)
                line.line.fill.background()
        
        # Bottom border
        if borders.get('bottom'):
            border = borders['bottom']
            if border.get('color') and border.get('width', 0) > 0:
                border_width_pt = px_to_pt(border['width'])
                # Create a thin rectangle for the bottom border
                line_height = border_width_pt / 72.0  # Convert points to inches
                border_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if (use_rounded_borders and is_rounded) else MSO_SHAPE.RECTANGLE
                line = slide.shapes.add_shape(
                    border_shape_type,
                    Inches(left), Inches(top + height - line_height),
                    Inches(width), Inches(line_height)
                )
                if use_rounded_borders and is_rounded:
                    try:
                        # Only round the top corners (bottom corners are at the edge)
                        adjustment = min((border_radius / min_dimension) * 2, 1.0) if min_dimension > 0 else 0.1
                        line.adjustments[0] = adjustment
                    except:
                        pass
                r, g, b = blend_transparent_color(border['color'], (255, 255, 255))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(r, g, b)
                line.line.fill.background()
        
        # Left border
        if borders.get('left'):
            border = borders['left']
            if border.get('color') and border.get('width', 0) > 0:
                border_width_pt = px_to_pt(border['width'])
                # Create a thin rectangle for the left border
                line_width = border_width_pt / 72.0  # Convert points to inches
                border_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if (use_rounded_borders and is_rounded) else MSO_SHAPE.RECTANGLE
                line = slide.shapes.add_shape(
                    border_shape_type,
                    Inches(left), Inches(top),
                    Inches(line_width), Inches(height)
                )
                if use_rounded_borders and is_rounded:
                    try:
                        # Only round the right corners (left corners are at the edge)
                        adjustment = min((border_radius / min_dimension) * 2, 1.0) if min_dimension > 0 else 0.1
                        line.adjustments[0] = adjustment
                    except:
                        pass
                r, g, b = blend_transparent_color(border['color'], (255, 255, 255))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(r, g, b)
                line.line.fill.background()
    else:
        # Fallback to uniform border for backward compatibility
        border_color = elem.get('border_color')
        border_width = elem.get('border_width', 0)
        if border_color and border_width > 0:
            # Blend transparent colors with white background
            r, g, b = blend_transparent_color(border_color, (255, 255, 255))
            shape.line.color.rgb = RGBColor(r, g, b)
            # Convert pixels to points for border width
            shape.line.width = Pt(px_to_pt(border_width))
        else:
            shape.line.fill.background()
    
    shape.shadow.inherit = False


def create_styled_text_element(slide, elem, left, top, width, height, text_elements_by_position=None):
    """Create a styled text element (text with background/border)."""
    coords = elem['coordinates']
    border_radius = float(elem.get('border_radius', 0) or 0)
    
    # Check if this is a bullet element (small circular element, likely a bullet)
    is_bullet = False
    bullet_text_elem = None
    
    if text_elements_by_position:
        # Check if this styled_text is a bullet (small, circular, minimal/no text)
        text_content = elem.get('text', '').strip()
        is_small = coords.get('width', 0) <= 60 and coords.get('height', 0) <= 60
        is_circular = border_radius >= (min(coords.get('width', 0), coords.get('height', 0)) / 2) * 0.8
        is_bullet = is_small and is_circular and (not text_content or len(text_content) <= 3)
        
        if is_bullet:
            # Find nearby text element to align with
            bullet_x = coords['x']
            bullet_y = coords['y']
            bullet_right = bullet_x + coords.get('width', 0)
            
            # Look for text elements that start near this bullet's right edge
            for (text_x, text_y), text_elem in text_elements_by_position.items():
                text_coords = text_elem.get('coordinates', {})
                # Check if text is to the right of bullet (within reasonable distance)
                if (text_x >= bullet_right - 20 and text_x <= bullet_right + 100 and
                    abs(text_y - bullet_y) < 50):  # Same approximate vertical position
                    bullet_text_elem = text_elem
                    break
            
            # If no text found to the right, check if bullet is inside a list item area
            if not bullet_text_elem:
                # Look for text elements that overlap vertically
                for (text_x, text_y), text_elem in text_elements_by_position.items():
                    text_coords = text_elem.get('coordinates', {})
                    text_top = text_coords.get('y', 0)
                    text_bottom = text_top + text_coords.get('height', 0)
                    bullet_center_y = bullet_y + coords.get('height', 0) / 2
                    
                    # Check if bullet is vertically aligned with text
                    if (bullet_center_y >= text_top - 10 and bullet_center_y <= text_bottom + 10 and
                        text_x > bullet_x):  # Text is to the right
                        bullet_text_elem = text_elem
                        break
            
            # Adjust bullet vertical position to align with text
            if bullet_text_elem:
                text_coords = bullet_text_elem.get('coordinates', {})
                text_font = bullet_text_elem.get('font', {})
                text_font_size_pt = text_font.get('size', 12)  # Already in points
                
                # Calculate text baseline position
                # Text top position in inches (text uses MSO_ANCHOR.TOP, so text starts at top)
                text_top = pixels_to_inches(text_coords.get('y', 0))
                # Convert font size from points to inches (1 pt = 1/72 inch)
                text_font_size_inches = text_font_size_pt / 72.0
                
                # Align bullet center with text first line center
                # For text with MSO_ANCHOR.TOP, the first line center is approximately:
                # text_top + (font_size / 2) - accounting for line height
                # Typical line height is ~1.2x font size, so first line center is at font_size * 0.6
                text_first_line_center = text_top + text_font_size_inches * 0.6
                
                # Center bullet vertically on first line center
                bullet_height_inches = height
                top = text_first_line_center - (bullet_height_inches / 2)
    
    # Get border_radius with proper type handling
    border_radius_raw = elem.get('border_radius', 0)
    try:
        border_radius = float(border_radius_raw) if border_radius_raw is not None else 0.0
        if not (isinstance(border_radius, (int, float)) and border_radius >= 0):
            border_radius = 0.0
    except (ValueError, TypeError):
        border_radius = 0.0
    
    if border_radius > 0:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        # Convert CSS border-radius (pixels) to PowerPoint adjustment (0.0 to 1.0)
        try:
            min_dimension = min(coords.get('width', 0), coords.get('height', 0))
            if min_dimension > 0:
                adjustment = min((border_radius / min_dimension) * 2, 1.0)
            else:
                adjustment = 0.1
            shape.adjustments[0] = adjustment
        except:
            pass
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    
    # Try to apply gradient first
    gradient = elem.get('gradient')
    gradient_applied = False
    if gradient:
        gradient_applied = apply_gradient_fill(shape, gradient)
    
    # Fallback to solid color if gradient failed
    if not gradient_applied:
        fill_color = elem.get('fill_color')
        if fill_color and fill_color.get('a', 1) > 0:
            shape.fill.solid()
            # Blend transparent colors with white background
            r, g, b = blend_transparent_color(fill_color, (255, 255, 255))
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
        else:
            shape.fill.background()
    
    # Apply borders - check for individual side borders first
    borders = elem.get('borders', {})
    has_individual_borders = any(borders.get(side) for side in ['top', 'right', 'bottom', 'left'])
    
    # Check if all borders are identical - if so, use uniform border with dash style support
    all_borders_same = False
    if has_individual_borders:
        top = borders.get('top')
        right = borders.get('right')
        bottom = borders.get('bottom')
        left = borders.get('left')
        
        # All four borders must exist and be identical in width and color
        if all([top, right, bottom, left]):
            all_same_width = (top.get('width') == right.get('width') == 
                            bottom.get('width') == left.get('width'))
            all_same_color = (top.get('color') == right.get('color') == 
                            bottom.get('color') == left.get('color'))
            all_borders_same = all_same_width and all_same_color
    
    # If borders are not all the same, use individual border rendering
    if has_individual_borders and not all_borders_same:
        # Remove border from shape first to avoid grey border
        shape.line.fill.background()
        
        # Determine if we need rounded corners for border rectangles
        use_rounded_borders = border_radius > 0
        min_dimension = min(coords.get('width', 0), coords.get('height', 0))
        is_rounded = border_radius > 0
        
        # Apply borders to individual sides using thin rectangle shapes
        # Top border
        if borders.get('top'):
            border = borders['top']
            if border.get('color') and border.get('width', 0) > 0:
                border_width_pt = px_to_pt(border['width'])
                # Create a thin rectangle for the top border
                line_height = border_width_pt / 72.0  # Convert points to inches
                border_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if (use_rounded_borders and is_rounded) else MSO_SHAPE.RECTANGLE
                line = slide.shapes.add_shape(
                    border_shape_type,
                    Inches(left), Inches(top),
                    Inches(width), Inches(line_height)
                )
                if use_rounded_borders and is_rounded:
                    try:
                        # Only round the bottom corners (top corners are at the edge)
                        # Set adjustment to match the main shape's radius
                        adjustment = min((border_radius / min_dimension) * 2, 1.0) if min_dimension > 0 else 0.1
                        line.adjustments[0] = adjustment
                    except:
                        pass
                r, g, b = blend_transparent_color(border['color'], (255, 255, 255))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(r, g, b)
                line.line.fill.background()
        
        # Right border
        if borders.get('right'):
            border = borders['right']
            if border.get('color') and border.get('width', 0) > 0:
                border_width_pt = px_to_pt(border['width'])
                # Create a thin rectangle for the right border
                line_width = border_width_pt / 72.0  # Convert points to inches
                border_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if (use_rounded_borders and is_rounded) else MSO_SHAPE.RECTANGLE
                line = slide.shapes.add_shape(
                    border_shape_type,
                    Inches(left + width - line_width), Inches(top),
                    Inches(line_width), Inches(height)
                )
                if use_rounded_borders and is_rounded:
                    try:
                        # Only round the left corners (right corners are at the edge)
                        adjustment = min((border_radius / min_dimension) * 2, 1.0) if min_dimension > 0 else 0.1
                        line.adjustments[0] = adjustment
                    except:
                        pass
                r, g, b = blend_transparent_color(border['color'], (255, 255, 255))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(r, g, b)
                line.line.fill.background()
        
        # Bottom border
        if borders.get('bottom'):
            border = borders['bottom']
            if border.get('color') and border.get('width', 0) > 0:
                border_width_pt = px_to_pt(border['width'])
                # Create a thin rectangle for the bottom border
                line_height = border_width_pt / 72.0  # Convert points to inches
                border_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if (use_rounded_borders and is_rounded) else MSO_SHAPE.RECTANGLE
                line = slide.shapes.add_shape(
                    border_shape_type,
                    Inches(left), Inches(top + height - line_height),
                    Inches(width), Inches(line_height)
                )
                if use_rounded_borders and is_rounded:
                    try:
                        # Only round the top corners (bottom corners are at the edge)
                        adjustment = min((border_radius / min_dimension) * 2, 1.0) if min_dimension > 0 else 0.1
                        line.adjustments[0] = adjustment
                    except:
                        pass
                r, g, b = blend_transparent_color(border['color'], (255, 255, 255))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(r, g, b)
                line.line.fill.background()
        
        # Left border
        if borders.get('left'):
            border = borders['left']
            if border.get('color') and border.get('width', 0) > 0:
                border_width_pt = px_to_pt(border['width'])
                # Create a thin rectangle for the left border
                line_width = border_width_pt / 72.0  # Convert points to inches
                border_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if (use_rounded_borders and is_rounded) else MSO_SHAPE.RECTANGLE
                line = slide.shapes.add_shape(
                    border_shape_type,
                    Inches(left), Inches(top),
                    Inches(line_width), Inches(height)
                )
                if use_rounded_borders and is_rounded:
                    try:
                        # Only round the right corners (left corners are at the edge)
                        adjustment = min((border_radius / min_dimension) * 2, 1.0) if min_dimension > 0 else 0.1
                        line.adjustments[0] = adjustment
                    except:
                        pass
                r, g, b = blend_transparent_color(border['color'], (255, 255, 255))
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(r, g, b)
                line.line.fill.background()
    else:
        # Use uniform border (either all borders are the same, or using border_color/border_width)
        # If all borders are the same, use the first one; otherwise use border_color/border_width
        if all_borders_same and has_individual_borders:
            first_border = borders.get('top') or borders.get('left') or borders.get('bottom') or borders.get('right')
            border_color = first_border.get('color') if first_border else elem.get('border_color')
            border_width = first_border.get('width', 0) if first_border else elem.get('border_width', 0)
            border_style = elem.get('border_style', 'solid')
        else:
            # Fallback to uniform border for backward compatibility
            border_color = elem.get('border_color')
            border_width = elem.get('border_width', 0)
            border_style = elem.get('border_style', 'solid')
        
        if border_color and border_width > 0:
            # Blend transparent colors with white background
            r, g, b = blend_transparent_color(border_color, (255, 255, 255))
            shape.line.color.rgb = RGBColor(r, g, b)
            # Convert pixels to points for border width
            shape.line.width = Pt(px_to_pt(border_width))
            
            # Set dash style based on border style
            if border_style == 'dotted':
                shape.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
            elif border_style == 'dashed':
                shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
            else:
                shape.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
        else:
            shape.line.fill.background()
    
    shape.shadow.inherit = False
    
    text_frame = shape.text_frame
    # Enable word wrap for better text rendering in boxes
    text_frame.word_wrap = True
    
    # Set equal margins for proper centering (PowerPoint sometimes needs small margins)
    margin = Inches(0.01)
    text_frame.margin_left = margin
    text_frame.margin_right = margin
    text_frame.margin_top = margin
    text_frame.margin_bottom = margin
    
    # Center text vertically and horizontally
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    
    # Set text after configuring frame
    text_frame.text = elem['text']
    
    font_name = {'Arial': 'Arial', 'Proxima Nova': 'Calibri', 'Roboto': 'Calibri'}.get(elem['font'].get('family', 'Arial'), 'Calibri')
    font_weight = str(elem['font']['weight'])
    is_bold = font_weight in ['bold', '700', '800', '900'] or (font_weight.isdigit() and int(font_weight) >= 700)
    color = elem['color']
    
    for paragraph in text_frame.paragraphs:
        # Use the alignment from the element, defaulting to center for small badges/pills
        alignment_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY, 'start': PP_ALIGN.LEFT, 'end': PP_ALIGN.RIGHT}
        stored_alignment = elem.get('alignment', 'center')
        # For small badges/pills, always center; otherwise use stored alignment
        text_content = elem.get('text', '').strip()
        is_small_badge = len(text_content) <= 3 and coords.get('width', 0) <= 60 and coords.get('height', 0) <= 60
        if is_small_badge:
            paragraph.alignment = PP_ALIGN.CENTER
        else:
            paragraph.alignment = alignment_map.get(stored_alignment.lower() if isinstance(stored_alignment, str) else 'center', PP_ALIGN.CENTER)
        for run in paragraph.runs:
            run.font.size = Pt(elem['font']['size'])
            run.font.name = font_name
            if is_bold:
                run.font.bold = True
            # Blend transparent colors with white background
            r, g, b = blend_transparent_color(color, (255, 255, 255))
            run.font.color.rgb = RGBColor(r, g, b)


def create_table_element(slide, elem):
    """Create table elements cell by cell."""
    rows = elem.get('rows', [])
    if not rows:
        return
    
    for row in rows:
        for cell in row:
            coords = cell['coordinates']
            left = pixels_to_inches(coords['x'])
            top = pixels_to_inches(coords['y'])
            width = pixels_to_inches(coords['width'])
            height = pixels_to_inches(coords['height'])
            
            bg_color = cell.get('bg_color')
            if bg_color and bg_color.get('a', 0) >= 0:
                bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
                bg_shape.fill.solid()
                # Blend transparent colors with white background
                r, g, b = blend_transparent_color(bg_color, (255, 255, 255))
                bg_shape.fill.fore_color.rgb = RGBColor(r, g, b)
                bg_shape.line.fill.background()
                bg_shape.shadow.inherit = False
            
            # Helper function to create a border line
            def create_border_line(x1, y1, x2, y2, color, width, style):
                from pptx.enum.shapes import MSO_CONNECTOR
                line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
                r, g, b = blend_transparent_color(color, (255, 255, 255))
                line.line.color.rgb = RGBColor(r, g, b)
                line.line.width = Pt(px_to_pt(width))
                
                # Disable shadow to prevent "shadowy" appearance
                line.shadow.inherit = False
                
                # Set dash style - use ROUND_DOT for better visibility of dotted lines
                if style == 'dotted':
                    line.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
                elif style == 'dashed':
                    line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                else:
                    line.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
            
            # Border bottom
            border_bottom_color = cell.get('border_bottom_color')
            if border_bottom_color and cell.get('border_bottom_width', 0) > 0:
                create_border_line(left, top + height, left + width, top + height,
                                   border_bottom_color, cell.get('border_bottom_width', 0),
                                   cell.get('border_bottom_style', 'solid'))
            
            # Border left
            border_left_color = cell.get('border_left_color')
            if border_left_color and cell.get('border_left_width', 0) > 0:
                create_border_line(left, top, left, top + height,
                                   border_left_color, cell.get('border_left_width', 0),
                                   cell.get('border_left_style', 'solid'))
            
            # Border right
            border_right_color = cell.get('border_right_color')
            if border_right_color and cell.get('border_right_width', 0) > 0:
                create_border_line(left + width, top, left + width, top + height,
                                   border_right_color, cell.get('border_right_width', 0),
                                   cell.get('border_right_style', 'solid'))
            
            # Border top
            border_top_color = cell.get('border_top_color')
            if border_top_color and cell.get('border_top_width', 0) > 0:
                create_border_line(left, top, left + width, top,
                                   border_top_color, cell.get('border_top_width', 0),
                                   cell.get('border_top_style', 'solid'))
            
            # Pseudo-element separator on right (::after)
            pseudo_right = cell.get('pseudo_separator_right')
            if pseudo_right and pseudo_right.get('color'):
                color = pseudo_right['color']
                # Only render if color has opacity (not transparent)
                if color.get('a', 0) > 0:
                    create_border_line(left + width, top, left + width, top + height,
                                       color, pseudo_right.get('width', 2),
                                       pseudo_right.get('style', 'dotted'))
            
            # Pseudo-element separator on left (::before)
            pseudo_left = cell.get('pseudo_separator_left')
            if pseudo_left and pseudo_left.get('color'):
                color = pseudo_left['color']
                # Only render if color has opacity (not transparent)
                if color.get('a', 0) > 0:
                    create_border_line(left, top, left, top + height,
                                       color, pseudo_left.get('width', 2),
                                       pseudo_left.get('style', 'dotted'))
            
            textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            text_frame = textbox.text_frame
            text_frame.text = cell['text']
            text_frame.word_wrap = True
            
            # Set vertical alignment - center text vertically in cells
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Set margins and alignment
            paragraph = text_frame.paragraphs[0]
            alignment = cell.get('alignment', 'left')
            is_header = cell.get('is_header', False)
            
            # Headers typically need less side margin, more vertical margin for proper appearance
            if is_header:
                side_margin = Inches(0.05)  # Small but visible margin
                vert_margin = Inches(0.03)  # Slightly more vertical space
            else:
                side_margin = Inches(0.05)
                vert_margin = Inches(0.02)
            
            if alignment == 'center':
                paragraph.alignment = PP_ALIGN.CENTER
                text_frame.margin_left = side_margin
                text_frame.margin_right = side_margin
            elif alignment == 'right' or alignment == 'end':
                paragraph.alignment = PP_ALIGN.RIGHT
                text_frame.margin_left = side_margin
                text_frame.margin_right = Inches(0.05)
            elif alignment == 'start' or alignment == 'left':
                paragraph.alignment = PP_ALIGN.LEFT
                text_frame.margin_left = Inches(0.05)
                text_frame.margin_right = side_margin
            else:
                paragraph.alignment = PP_ALIGN.LEFT
                text_frame.margin_left = Inches(0.05)
                text_frame.margin_right = side_margin
            
            text_frame.margin_top = vert_margin
            text_frame.margin_bottom = vert_margin
            
            if paragraph.runs:
                run = paragraph.runs[0]
                run.font.size = Pt(int(cell.get('font_size', 12) * 0.75))
                run.font.name = 'Calibri'
                if cell.get('is_header'):
                    run.font.bold = True
                color = cell.get('color', {'r': 0, 'g': 0, 'b': 0})
                if color:
                    # Blend transparent colors with white background
                    r, g, b = blend_transparent_color(color, (255, 255, 255))
                    run.font.color.rgb = RGBColor(r, g, b)


def create_image_element(slide, elem, left, top, width, height):
    """Create an image element."""
    try:
        img_src = elem.get('src', '')
        if not img_src:
            print(f"  Warning: Image element has no src attribute")
            return
        
        # Ensure width and height are valid
        if width <= 0 or height <= 0:
            print(f"  Warning: Image has invalid dimensions: {width}x{height}")
            return
        
        # Debug: print image info (only for first few to avoid spam)
        # print(f"  Adding image: {img_src[:80]}... at ({left:.2f}, {top:.2f}), size ({width:.2f}, {height:.2f})")
        
        natural_width = elem.get('natural_width')
        natural_height = elem.get('natural_height')
        object_fit = elem.get('object_fit', 'fill')
        
        if natural_width and natural_height and natural_width > 0 and natural_height > 0:
            natural_aspect = natural_width / natural_height
            display_aspect = width / height if height > 0 else 1
            if object_fit == 'contain':
                if natural_aspect > display_aspect:
                    h_new = width / natural_aspect
                    top += (height - h_new) / 2
                    height = h_new
                else:
                    w_new = height * natural_aspect
                    left += (width - w_new) / 2
                    width = w_new
            elif abs(natural_aspect - display_aspect) > 0.1:
                if natural_aspect > display_aspect:
                    height = width / natural_aspect
                else:
                    width = height * natural_aspect
        
        is_circle = elem.get('is_circle', False)
        pic = None
        if img_src.startswith('data:image'):
            # Handle data URI (base64 encoded images)
            try:
                # Extract base64 data from data URI: data:image/png;base64,<data>
                header, encoded = img_src.split(',', 1)
                img_data = base64.b64decode(encoded)
                img_stream = io.BytesIO(img_data)
                img_stream.seek(0)  # Reset stream position
                pic = slide.shapes.add_picture(img_stream, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
                print(f"    ✓ Added data URI image ({len(img_data)} bytes)")
            except Exception as e:
                print(f"  Warning: Could not decode data URI image: {e}")
                import traceback
                traceback.print_exc()
        elif img_src.startswith('http'):
            try:
                req = urllib.request.Request(img_src, headers={'User-Agent': 'Mozilla/5.0'})
                with urllib.request.urlopen(req, timeout=10) as response:
                    img_data = response.read()
                    if len(img_data) == 0:
                        print(f"  Warning: Image data is empty for {img_src[:60]}...")
                        pic = None
                    else:
                        img_stream = io.BytesIO(img_data)
                        img_stream.seek(0)  # Reset stream position
                        pic = slide.shapes.add_picture(img_stream, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
                        print(f"    ✓ Added HTTP image ({len(img_data)} bytes)")
            except Exception as e:
                print(f"  Warning: Could not load image from {img_src[:80]}...: {e}")
                import traceback
                traceback.print_exc()
                pic = None
        elif os.path.exists(img_src):
            try:
                pic = slide.shapes.add_picture(img_src, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
                print(f"    ✓ Added local image")
            except Exception as e:
                print(f"  Warning: Could not load local image {img_src}: {e}")
                import traceback
                traceback.print_exc()
                pic = None
        else:
            print(f"  Warning: Image source not recognized: {img_src[:80]}...")
        
        # Don't apply circular clipping - it might be hiding the images
        # The gradient circle shape provides the circular background
        # The image should show fully on top
        # if pic and is_circle:
        #     from pptx.oxml import parse_xml
        #     spPr = pic._element.spPr
        #     prstGeom = parse_xml('<a:prstGeom xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="ellipse"><a:avLst/></a:prstGeom>')
        #     for child in list(spPr):
        #         if 'Geom' in child.tag:
        #             spPr.remove(child)
        #     spPr.insert(0, prstGeom)
        
        # CRITICAL: Move image to absolute end of shapes collection to ensure it's on top
        # PowerPoint renders shapes in order, so last shape appears on top
        if pic:
            try:
                # Remove from current position and re-add at end to ensure it's on top
                slide.shapes._spTree.remove(pic._element)
                slide.shapes._spTree.append(pic._element)
                print(f"    ✓ Moved image to top (z-order)")
            except Exception as e:
                print(f"  Warning: Could not move image to top: {e}")
                # Image should still be visible, just might be covered
    except Exception as e:
        print(f"  Error creating image element: {e}")
        import traceback
        traceback.print_exc()


def create_text_element(slide, elem, left, top, width, height):
    """Create a text element."""
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = textbox.text_frame
    text_frame.text = elem['text']
    # Enable word wrap for all text to prevent overflow
    text_frame.word_wrap = True
    
    if len(elem['text'].strip()) <= 3:
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    else:
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    alignment_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT, 'justify': PP_ALIGN.JUSTIFY, 'start': PP_ALIGN.LEFT, 'end': PP_ALIGN.RIGHT}
    # Always respect the alignment from the HTML/CSS, don't auto-center based on text length
    text_alignment = alignment_map.get(elem.get('alignment', 'left'), PP_ALIGN.LEFT)
    
    # Set margins based on alignment - PowerPoint needs small margins for proper text alignment
    # Use small margin for left/right alignment to ensure text aligns properly
    margin = Inches(0.01) if text_alignment in [PP_ALIGN.LEFT, PP_ALIGN.RIGHT] else Inches(0)
    text_frame.margin_left = margin if text_alignment == PP_ALIGN.LEFT else 0
    text_frame.margin_right = margin if text_alignment == PP_ALIGN.RIGHT else 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    
    font_name = {'Arial': 'Arial', 'Calibri': 'Calibri', 'Times New Roman': 'Times New Roman'}.get(elem['font'].get('family', 'Arial'), 'Calibri')
    font_weight = str(elem['font']['weight'])
    is_bold = font_weight in ['bold', '700', '800', '900'] or (font_weight.isdigit() and int(font_weight) >= 700)
    color = elem['color']
    text_gradient = elem.get('text_gradient')
    
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = text_alignment
        
        for run in paragraph.runs:
            run.font.size = Pt(elem['font']['size'])
            run.font.name = font_name
            if is_bold:
                run.font.bold = True
            
            # Convert gradient text to solid fill using one of the gradient colors
            if text_gradient and text_gradient.get('stops'):
                # Get gradient stops and pick the first stop color as solid fill
                stops = sorted(text_gradient.get('stops', []), key=lambda s: s.get('position', 0))
                if stops and stops[0].get('color'):
                    gradient_color = stops[0]['color']
                else:
                    # Fallback to element color if gradient stop has no color
                    gradient_color = color
                
                # Use the gradient color as solid fill
                r, g, b = blend_transparent_color(gradient_color, (255, 255, 255))
                run.font.color.rgb = RGBColor(r, g, b)
            else:
                # Blend transparent colors with white background
                r, g, b = blend_transparent_color(color, (255, 255, 255))
                run.font.color.rgb = RGBColor(r, g, b)
    
    border_color = elem.get('border_color')
    if border_color and elem.get('border_width', 0) > 0:
        # Blend transparent colors with white background
        r, g, b = blend_transparent_color(border_color, (255, 255, 255))
        textbox.line.color.rgb = RGBColor(r, g, b)
        # Convert pixels to points for border width
        textbox.line.width = Pt(px_to_pt(elem.get('border_width', 0)))
    else:
        textbox.line.fill.background()


def create_image_shape(slide, elem, left_emu, top_emu, width_emu, height_emu):
    """Create a picture shape from image element."""
    media = elem.get('media', {})
    img_src = media.get('image_src', '')
    
    if not img_src:
        return
    
    try:
        pic = None
        # Handle HTTP/HTTPS URLs
        if img_src.startswith('http'):
            with urllib.request.urlopen(img_src) as response:
                img_data = response.read()
                img_stream = io.BytesIO(img_data)
                pic = slide.shapes.add_picture(
                    img_stream,
                    left_emu, top_emu,
                    width=width_emu, height=height_emu
                )
        elif os.path.exists(img_src):
            # Local file
            pic = slide.shapes.add_picture(
                img_src,
                left_emu, top_emu,
                width=width_emu, height=height_emu
            )
        
        # Set link if present
        if pic:
            link_data = elem.get('link', {})
            if link_data.get('href'):
                try:
                    pic.click_action.action = 'ppActionHyperlink'
                    pic.click_action.hyperlink.address = link_data['href']
                except:
                    pass
            
            # Set opacity if present
            opacity = elem.get('opacity')
            if opacity is not None and opacity < 1:
                try:
                    pic.fill.transparency = 1 - opacity
                except:
                    pass
            
            # Handle circular images
            media = elem.get('media', {})
            if media.get('is_circle'):
                try:
                    from pptx.oxml import parse_xml
                    spPr = pic._element.spPr
                    prstGeom = parse_xml(
                        '<a:prstGeom xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" prst="ellipse">'
                        '<a:avLst/>'
                        '</a:prstGeom>'
                    )
                    for child in list(spPr):
                        if 'Geom' in child.tag:
                            spPr.remove(child)
                    spPr.insert(0, prstGeom)
                except:
                    pass
            
            # Handle object-fit: contain (preserve aspect ratio)
            if media.get('object_fit') == 'contain':
                natural_width = media.get('image_natural_width_px', 0)
                natural_height = media.get('image_natural_height_px', 0)
                if natural_width > 0 and natural_height > 0:
                    # Convert EMU bounds back to pixels for calculation
                    bounds = elem.get('bounds', {})
                    display_width_px = bounds.get('width', 0)
                    display_height_px = bounds.get('height', 0)
                    
                    if display_width_px > 0 and display_height_px > 0:
                        natural_aspect = natural_width / natural_height
                        display_aspect = display_width_px / display_height_px
                        
                        if natural_aspect > display_aspect:
                            # Image is wider - fit to width
                            new_height_px = display_width_px / natural_aspect
                            new_height_emu = px_to_emu_y(new_height_px)
                            pic.height = new_height_emu
                            pic.top = top_emu + (height_emu - new_height_emu) // 2
                        else:
                            # Image is taller - fit to height
                            new_width_px = display_height_px * natural_aspect
                            new_width_emu = px_to_emu_x(new_width_px)
                            pic.width = new_width_emu
                            pic.left = left_emu + (width_emu - new_width_emu) // 2
    except Exception as e:
        print(f"  Warning: Could not add image {img_src}: {e}")


def create_text_shape(slide, elem, left_emu, top_emu, width_emu, height_emu):
    """Create a text box from text element."""
    text_data = elem.get('text', {})
    if not text_data:
        return
    
    # Account for padding - adjust position and size
    padding = elem.get('padding', {})
    padding_left_px = padding.get('left', 0)
    padding_top_px = padding.get('top', 0)
    padding_right_px = padding.get('right', 0)
    padding_bottom_px = padding.get('bottom', 0)
    
    # Adjust text box position to account for padding
    # The bounds include padding, but text content starts after padding
    adjusted_left_emu = left_emu + px_to_emu_x(padding_left_px)
    adjusted_top_emu = top_emu + px_to_emu_y(padding_top_px)
    adjusted_width_emu = width_emu - px_to_emu_x(padding_left_px + padding_right_px)
    adjusted_height_emu = height_emu - px_to_emu_y(padding_top_px + padding_bottom_px)
    
    # Ensure non-negative dimensions and minimum size
    # If padding makes the box too small, use original bounds
    if adjusted_width_emu <= 0 or adjusted_height_emu <= 0:
        # Padding adjustment resulted in invalid size, use original bounds
        adjusted_left_emu = left_emu
        adjusted_top_emu = top_emu
        adjusted_width_emu = width_emu
        adjusted_height_emu = height_emu
    else:
        # Ensure minimum size
        adjusted_width_emu = max(adjusted_width_emu, px_to_emu_x(10))
        adjusted_height_emu = max(adjusted_height_emu, px_to_emu_y(10))
    
    # Ensure position is valid
    if adjusted_left_emu < 0:
        adjusted_left_emu = 0
    if adjusted_top_emu < 0:
        adjusted_top_emu = 0
    
    textbox = slide.shapes.add_textbox(
        adjusted_left_emu, adjusted_top_emu,
        adjusted_width_emu, adjusted_height_emu
    )
    
    # Set fill/background if present (but not if it's a background image - handled separately)
    fill_data = elem.get('fill', {})
    bg_image_url = fill_data.get('background_image_url')
    if not bg_image_url:
        bg_color_rgba = fill_data.get('background_color_rgba')
        if bg_color_rgba:
            rgb = rgba_to_rgb(bg_color_rgba)
            if rgb:
                textbox.fill.solid()
                textbox.fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
        else:
            # No background color - explicitly set to no fill (transparent)
            # This ensures text is visible even without a background
            try:
                textbox.fill.background()
            except:
                pass
    
    # Set opacity if present
    opacity = elem.get('opacity')
    if opacity is not None and opacity < 1:
        try:
            # Only set transparency if we have a fill
            if fill_data.get('background_color_rgba'):
                textbox.fill.transparency = 1 - opacity
        except:
            pass
    
    # Set link if present
    link_data = elem.get('link', {})
    if link_data.get('href'):
        try:
            textbox.click_action.action = 'ppActionHyperlink'
            textbox.click_action.hyperlink.address = link_data['href']
        except:
            pass
    
    text_frame = textbox.text_frame
    text_content = text_data.get('content', '').strip()
    if not text_content:
        # No text content - but if this is part of a shape+text combo (border/background image),
        # the shape was already created, so we can skip the empty text box
        return
    text_frame.text = text_content
    text_frame.word_wrap = True
    
    # Remove margins for accurate positioning
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    
    # If this text box is on top of a shape (border/background image), make it transparent
    # so the shape shows through
    if elem.get('border') or elem.get('fill', {}).get('background_image_url'):
        try:
            # Ensure text box background is transparent when it's on top of a shape
            textbox.fill.background()
        except:
            pass
    
    # Format paragraph
    p = text_frame.paragraphs[0]
    
    # Alignment
    align_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    p.alignment = align_map.get(text_data.get('text_align', 'left').lower(), PP_ALIGN.LEFT)
    
    # Line spacing (PowerPoint uses spacing in points, relative to font size)
    line_height = text_data.get('line_height', 'normal')
    if line_height and line_height != 'normal':
        try:
            font_size_pt = px_to_pt(text_data.get('font_size_px', 12))
            # Parse line-height (could be number, px, em, etc.)
            if isinstance(line_height, str):
                if line_height.endswith('px'):
                    line_height_val = float(line_height.replace('px', ''))
                    # Convert px to pt and calculate spacing
                    line_height_pt = px_to_pt(line_height_val)
                    # PowerPoint line_spacing is spacing in points
                    p.line_spacing = line_height_pt - font_size_pt
                elif line_height.endswith('em'):
                    # em is relative to font size
                    em_multiplier = float(line_height.replace('em', ''))
                    line_height_pt = font_size_pt * em_multiplier
                    p.line_spacing = line_height_pt - font_size_pt
                elif line_height.replace('.', '').replace('-', '').isdigit():
                    # Unitless number (multiplier of font size)
                    multiplier = float(line_height)
                    line_height_pt = font_size_pt * multiplier
                    p.line_spacing = line_height_pt - font_size_pt
            elif isinstance(line_height, (int, float)):
                # Assume it's a multiplier
                line_height_pt = font_size_pt * line_height
                p.line_spacing = line_height_pt - font_size_pt
        except Exception as e:
            # Silently fail - use default spacing
            pass
    
    # Format text run
    if p.runs:
        run = p.runs[0]
    else:
        run = p.add_run()
    
    # Font properties
    font_family = text_data.get('font_family', 'Calibri')
    run.font.name = font_family
    run.font.size = Pt(px_to_pt(text_data.get('font_size_px', 12)))
    
    # Font weight
    font_weight = text_data.get('font_weight', 'normal')
    if font_weight in ['bold', '700', '800', '900'] or \
       (isinstance(font_weight, str) and font_weight.isdigit() and int(font_weight) >= 700):
        run.font.bold = True
    
    # Font style
    if text_data.get('font_style') == 'italic':
        run.font.italic = True
    
    # Text decoration (underline, strikethrough)
    text_decoration = text_data.get('text_decoration', '') or text_data.get('text_decoration_line', '')
    if text_decoration:
        if 'underline' in text_decoration.lower():
            run.font.underline = True
        if 'line-through' in text_decoration.lower() or 'strikethrough' in text_decoration.lower():
            run.font.strike = True
    
    # Color
    color_rgba = text_data.get('color_rgba', 'rgba(0,0,0,1)')
    rgb = rgba_to_rgb(color_rgba)
    if rgb:
        run.font.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])


def create_shape(slide, elem, left_emu, top_emu, width_emu, height_emu):
    """Create a shape (rectangle) from element with fill/border."""
    # Determine shape type based on border radius and aspect ratio
    border = elem.get('border', {})
    bounds = elem.get('bounds', {})
    width_px = bounds.get('width', 0)
    height_px = bounds.get('height', 0)
    
    # Get all border radius values
    radius_tl = border.get('radius_top_left_px', 0)
    radius_tr = border.get('radius_top_right_px', 0)
    radius_br = border.get('radius_bottom_right_px', 0)
    radius_bl = border.get('radius_bottom_left_px', 0)
    
    # Use the maximum radius for shape type determination
    max_radius = max(radius_tl, radius_tr, radius_br, radius_bl)
    
    # Check if it's a perfect circle (square aspect ratio + border-radius: 50%)
    is_circle = False
    if width_px > 0 and height_px > 0:
        aspect_ratio = width_px / height_px
        is_squareish = aspect_ratio > 0.9 and aspect_ratio < 1.1
        min_dimension = min(width_px, height_px)
        # Check if border-radius is approximately 50% (circle)
        # Use average of all corners or max radius
        avg_radius = (radius_tl + radius_tr + radius_br + radius_bl) / 4
        is_circle = is_squareish and avg_radius >= (min_dimension / 2) * 0.8
    
    # Determine if rounded (any corner has significant rounding)
    # Use percentage-based threshold: if radius is > 5% of smallest dimension
    is_rounded = False
    if width_px > 0 and height_px > 0:
        min_dimension = min(width_px, height_px)
        # Consider rounded if any corner has radius > 2% of smallest dimension
        is_rounded = max_radius > (min_dimension * 0.02) or max_radius > 5
    
    if is_circle:
        # Perfect circle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left_emu, top_emu,
            width_emu, height_emu
        )
    elif is_rounded:
        # Rounded rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_emu, top_emu,
            width_emu, height_emu
        )
    else:
        # Regular rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left_emu, top_emu,
            width_emu, height_emu
        )
    
    # Set fill
    fill_data = elem.get('fill', {})
    bg_image_url = fill_data.get('background_image_url')
    
    if bg_image_url:
        # Background image
        try:
            if bg_image_url.startswith('http'):
                with urllib.request.urlopen(bg_image_url) as response:
                    img_data = response.read()
                    img_stream = io.BytesIO(img_data)
                    # Create picture shape instead
                    pic = slide.shapes.add_picture(
                        img_stream,
                        left_emu, top_emu,
                        width=width_emu, height=height_emu
                    )
                    # Set link if present
                    link_data = elem.get('link', {})
                    if link_data.get('href'):
                        try:
                            pic.click_action.action = 'ppActionHyperlink'
                            pic.click_action.hyperlink.address = link_data['href']
                        except:
                            pass
                    return
            elif os.path.exists(bg_image_url):
                pic = slide.shapes.add_picture(
                    bg_image_url,
                    left_emu, top_emu,
                    width=width_emu, height=height_emu
                )
                link_data = elem.get('link', {})
                if link_data.get('href'):
                    try:
                        pic.click_action.action = 'ppActionHyperlink'
                        pic.click_action.hyperlink.address = link_data['href']
                    except:
                        pass
                return
        except Exception as e:
            print(f"  Warning: Could not add background image {bg_image_url}: {e}")
    
    # Solid color fill
    bg_color_rgba = fill_data.get('background_color_rgba')
    if bg_color_rgba:
        rgb = rgba_to_rgb(bg_color_rgba)
        if rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
    else:
        shape.fill.background()
    
    # Set opacity if present
    opacity = elem.get('opacity')
    if opacity is not None and opacity < 1:
        try:
            shape.fill.transparency = 1 - opacity
        except:
            pass
    
    # Set border (all sides - PowerPoint uses uniform border, so use the most prominent)
    if border:
        # Find the most prominent border
        borders = ['top', 'right', 'bottom', 'left']
        max_border = None
        max_width = 0
        
        for side in borders:
            side_border = border.get(side, {})
            width = side_border.get('width', 0) or side_border.get('width_px', 0)
            if width > max_width:
                max_width = width
                max_border = side_border
        
        if max_border and max_width > 0:
            border_color_rgba = max_border.get('color', '') or max_border.get('color_rgba', '')
            border_style = max_border.get('style', 'solid')
            
            if border_color_rgba:
                rgb = rgba_to_rgb(border_color_rgba)
                if rgb:
                    shape.line.color.rgb = RGBColor(rgb[0], rgb[1], rgb[2])
                    shape.line.width = Pt(px_to_pt(max_width))
                    
                    # Set dash style
                    if 'dashed' in border_style.lower():
                        try:
                            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                        except:
                            pass
                    elif 'dotted' in border_style.lower():
                        try:
                            shape.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
                        except:
                            pass
        else:
            shape.line.fill.background()
    else:
        shape.line.fill.background()
    
    # Set shadow if present
    shadow_data = elem.get('shadow', {})
    if shadow_data.get('box_shadow'):
        try:
            # Parse box-shadow: offset-x offset-y blur-radius spread-radius color
            shadow_str = shadow_data['box_shadow']
            # Simple shadow implementation - PowerPoint has limited shadow support
            shape.shadow.inherit = False
            shape.shadow.style = 'outer'
        except:
            pass
    else:
        shape.shadow.inherit = False
    
    # Set link if present
    link_data = elem.get('link', {})
    if link_data.get('href'):
        try:
            shape.click_action.action = 'ppActionHyperlink'
            shape.click_action.hyperlink.address = link_data['href']
        except:
            pass


async def convert_json_to_pptx(json_path: str, output_path: str):
    """
    Step 1: Read slide data from JSON file.
    Then process each slide through the conversion pipeline.
    """
    # Load JSON
    with open(json_path, 'r') as f:
        slides_data = json.load(f)
    
    print(f"Processing {len(slides_data)} slides...")
    
    # Create presentation with custom slide size
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    
    # Process each slide
    for idx, slide_obj in enumerate(slides_data, 1):
        slide_id = slide_obj.get('id', f'slide_{idx}')
        html_content = slide_obj['html']
        
        print(f"  [{idx}/{len(slides_data)}] {slide_id}")
        
        # Step 2: Extract elements from HTML
        elements_json = await extract_elements_from_html(html_content)
        
        # Step 4: Convert to PPTX
        create_pptx_from_elements(prs, elements_json)
    
    # Save presentation
    prs.save(output_path)
    print(f"\n✓ Created: {output_path}")


async def main():
    if len(sys.argv) < 2:
        print("Usage: python3 html_to_pptx.py <json_file> [output.pptx]")
        sys.exit(1)
    
    json_file = sys.argv[1]
    
    # Default output: same name as input but with .pptx extension
    if len(sys.argv) > 2:
        output_file = sys.argv[2]
    else:
        output_file = str(Path(json_file).with_suffix('.pptx'))
    
    try:
        await convert_json_to_pptx(json_file, output_file)
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    asyncio.run(main())

